#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
楽々JC v2.0.0
─────────────────────────────────────────
  機能①  リンク一括設定（旧 LinkForge）
  機能②  PDF一括変換
─────────────────────────────────────────
Mac / Windows 対応
tkinterdnd2 によるドラッグ&ドロップ対応
"""

import os
import sys
import re
import csv
import shutil
import threading
import platform
import subprocess
import urllib.request
import urllib.error
import urllib.parse
from pathlib import Path
from copy import deepcopy
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

# ── python-docx ─────────────────────────────────────────────────
try:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.shared import RGBColor as DocxRGBColor, Pt as DocxPt
except ImportError:
    root = tk.Tk(); root.withdraw()
    messagebox.showerror("エラー",
        "python-docx がインストールされていません。\n\n"
        "ターミナルで以下を実行してください:\n\n"
        "【Mac】  python3 -m pip install python-docx\n"
        "【Win】  python -m pip install python-docx")
    sys.exit(1)

# ── openpyxl (チェッカー Excel処理) ──────────────────────────────
try:
    import openpyxl
    from openpyxl.styles import Font as OpenpyxlFont
    _OPENPYXL_OK = True
except ImportError:
    _OPENPYXL_OK = False

# ── python-pptx (チェッカー PPT処理) ─────────────────────────────
try:
    from pptx import Presentation as PptxPresentation
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.util import Pt as PptxPt
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

# ── pdfplumber (チェッカー PDF解析) ──────────────────────────────
try:
    import pdfplumber
    _PDFPLUMBER_OK = True
except ImportError:
    _PDFPLUMBER_OK = False

# ── D&D ──────────────────────────────────────────────────────────
_DND_BACKEND = None
try:
    from tkinterdnd2 import DND_FILES, TkinterDnD
    _DND_BACKEND = "tkdnd"
except ImportError:
    DND_FILES = "DND_Files"

# ════════════════════════════════════════════════════════════════
#  カラーテーマ（全画面共通）
# ════════════════════════════════════════════════════════════════
C = {
    "bg":       "#1B2E4A",
    "surface":  "#1E3759",
    "accent":   "#2A4F8A",
    "primary":  "#E94560",
    "text":     "#FFFFFF",
    "sub":      "#C8DCF5",
    "ok":       "#00D9A5",
    "warn":     "#FFC857",
    "err":      "#FF6B6B",
    "info":     "#48C9F7",
    "input_bg": "#152840",
    "border":   "#3A5A8A",
    "drop_hi":  "#274470",
}

APP_VERSION  = "2.0.0"
GITHUB_USER  = "domechin6248"
GITHUB_REPO  = "LinkForge"
GITHUB_RAW   = f"https://raw.githubusercontent.com/{GITHUB_USER}/{GITHUB_REPO}/main"

HYPERLINK_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/"
    "relationships/hyperlink"
)

LINK_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls",
    ".pptx", ".ppt", ".csv", ".txt",
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff",
    ".zip", ".rtf", ".odt", ".ods", ".odp",
}

PDF_EXTENSIONS = {
    ".doc", ".docx", ".odt", ".rtf",
    ".xls", ".xlsx", ".ods", ".csv",
    ".ppt", ".pptx", ".odp",
    ".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".gif", ".webp",
    ".txt",
}

# 変換エンジン振り分け用
_WORD_EXTS  = {".doc", ".docx", ".odt", ".rtf", ".txt"}
_EXCEL_EXTS = {".xls", ".xlsx", ".ods", ".csv"}
_PPT_EXTS   = {".ppt", ".pptx", ".odp"}
_IMG_EXTS   = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".gif", ".webp"}


# ════════════════════════════════════════════════════════════════
#  共通ユーティリティ
# ════════════════════════════════════════════════════════════════

def make_scrollable_frame(parent):
    """スクロール可能なフレームを返す (canvas, inner_frame)"""
    outer = tk.Frame(parent, bg=C["bg"])
    outer.pack(fill=tk.BOTH, expand=True)
    cv = tk.Canvas(outer, bg=C["bg"], highlightthickness=0, bd=0)
    sb = ttk.Scrollbar(outer, orient="vertical", command=cv.yview)
    cv.configure(yscrollcommand=sb.set)
    sb.pack(side=tk.RIGHT, fill=tk.Y)
    cv.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
    inner = tk.Frame(cv, bg=C["bg"])
    fid = cv.create_window((0, 0), window=inner, anchor="nw")
    inner.bind("<Configure>", lambda e: cv.configure(scrollregion=cv.bbox("all")))
    cv.bind("<Configure>", lambda e: cv.itemconfig(fid, width=e.width))
    return cv, inner


def make_log_widget(parent):
    """ログテキストウィジェットを返す"""
    log_wrap = tk.Frame(parent, bg=C["accent"], padx=1, pady=1)
    log_wrap.pack(fill=tk.BOTH, expand=True, padx=16, pady=(0, 16))
    tk.Label(log_wrap, text=" 処理ログ",
             font=("Helvetica", 9), bg=C["accent"], fg="#C8DCF5"
             ).pack(anchor="w")
    log_inner = tk.Frame(log_wrap, bg=C["input_bg"])
    log_inner.pack(fill=tk.BOTH, expand=True)
    mono = "Menlo" if platform.system() == "Darwin" else "Consolas"
    txt = tk.Text(
        log_inner, height=8, font=(mono, 9),
        bg=C["input_bg"], fg=C["sub"],
        insertbackground=C["primary"],
        relief=tk.FLAT, bd=0, padx=10, pady=8,
        state=tk.DISABLED, wrap=tk.WORD
    )
    lsb = ttk.Scrollbar(log_inner, command=txt.yview)
    txt.configure(yscrollcommand=lsb.set)
    lsb.pack(side=tk.RIGHT, fill=tk.Y)
    txt.pack(fill=tk.BOTH, expand=True)
    txt.tag_configure("success", foreground=C["ok"])
    txt.tag_configure("error",   foreground=C["err"])
    txt.tag_configure("info",    foreground=C["info"])
    txt.tag_configure("warning", foreground=C["warn"])
    return txt


def log_write(txt_widget, msg, tag=""):
    txt_widget.configure(state=tk.NORMAL)
    if tag:
        txt_widget.insert(tk.END, msg + "\n", tag)
    else:
        txt_widget.insert(tk.END, msg + "\n")
    txt_widget.see(tk.END)
    txt_widget.configure(state=tk.DISABLED)


def section_divider(parent):
    tk.Frame(parent, bg=C["border"], height=1).pack(fill=tk.X, padx=16, pady=(8, 10))


# ════════════════════════════════════════════════════════════════
#  FlatButton  ─  macOS でも bg/fg が確実に反映されるボタン
# ════════════════════════════════════════════════════════════════

class FlatButton(tk.Frame):
    """tk.Frame + tk.Label の組み合わせで macOS ネイティブテーマを回避。
    tk.Button と同じ引数でほぼ互換の configure/config を持つ。"""

    _IGNORE = frozenset({
        "activebackground", "activeforeground",
        "highlightbackground", "highlightthickness",
        "relief", "bd",
    })

    def __init__(self, parent, text="", command=None, bg=None, fg="#FFFFFF",
                 font=("Helvetica", 9), bold=False, padx=12, pady=4,
                 cursor="hand2", state=tk.NORMAL, **kw):
        for k in list(kw):
            if k in self._IGNORE:
                del kw[k]
        self._cmd    = command
        self._state  = state
        self._bg     = bg or C["accent"]
        self._fg     = fg
        self._cursor = cursor
        if bold:
            if isinstance(font, tuple) and len(font) >= 2:
                font = (font[0], font[1], "bold")
            else:
                font = ("Helvetica", 9, "bold")
        _cur = cursor if state == tk.NORMAL else "arrow"
        super().__init__(parent, bg=self._bg, cursor=_cur, **kw)
        self._lbl = tk.Label(
            self, text=text, bg=self._bg, fg=self._fg,
            font=font, padx=padx, pady=pady, cursor=_cur
        )
        self._lbl.pack(fill=tk.BOTH, expand=True)
        for w in (self, self._lbl):
            w.bind("<Button-1>", self._on_click)
            w.bind("<Enter>",    self._on_enter)
            w.bind("<Leave>",    self._on_leave)

    def _on_click(self, e=None):
        if self._state == tk.NORMAL and self._cmd:
            self._cmd()

    def _on_enter(self, e=None):
        if self._state == tk.NORMAL:
            super().configure(bg=C["primary"])
            self._lbl.configure(bg=C["primary"])

    def _on_leave(self, e=None):
        if self._state == tk.NORMAL:
            super().configure(bg=self._bg)
            self._lbl.configure(bg=self._bg)

    def configure(self, **kw):
        for k in list(kw):
            if k in self._IGNORE:
                del kw[k]
        text    = kw.pop("text",    None)
        state   = kw.pop("state",   None)
        bg      = kw.pop("bg",      None)
        fg      = kw.pop("fg",      None)
        cursor  = kw.pop("cursor",  None)
        command = kw.pop("command", None)
        if kw:
            super().configure(**kw)
        if command is not None:
            self._cmd = command
        if text is not None:
            self._lbl.configure(text=text)
        if bg is not None:
            self._bg = bg
        if fg is not None:
            self._fg = fg
        if cursor is not None:
            self._cursor = cursor
        if state is not None:
            self._state = state
        if self._state == tk.NORMAL:
            super().configure(bg=self._bg, cursor=self._cursor)
            self._lbl.configure(bg=self._bg, fg=self._fg, cursor=self._cursor)
        else:
            super().configure(bg=self._bg, cursor="arrow")
            self._lbl.configure(bg=self._bg, fg=self._fg, cursor="arrow")

    config = configure


# ════════════════════════════════════════════════════════════════
#  共通ボタンヘルパー
# ════════════════════════════════════════════════════════════════

def nav_button(parent, text, command):
    return FlatButton(
        parent, text=text, command=command,
        font=("Helvetica", 9), bold=True,
        bg=C["accent"], fg="#FFFFFF",
        padx=14, pady=5
    )


def make_btn(parent, text, command, font_size=9, bold=False,
             bg=None, fg="#FFFFFF", padx=12, pady=4, cursor="hand2", **kw):
    return FlatButton(
        parent, text=text, command=command,
        font=("Helvetica", font_size), bold=bold,
        bg=bg or C["accent"], fg=fg,
        padx=padx, pady=pady, cursor=cursor, **kw
    )


class LoggedFrame(tk.Frame):
    """ログウィジェット付き Frame の基底クラス。"""
    def _log(self, msg, tag=""):
        self.after(0, lambda m=msg, t=tag: log_write(self.log_txt, m, t))


# ════════════════════════════════════════════════════════════════
#  DropZone ウィジェット（共通）
# ════════════════════════════════════════════════════════════════

class DropZone(tk.Frame):

    def __init__(self, parent, label_text, hint_text,
                 select_mode="file", file_types=None,
                 allow_multiple=False, count_extensions=None, **kwargs):
        super().__init__(parent, bg=C["surface"], **kwargs)
        self.select_mode      = select_mode
        self.file_types       = file_types or []
        self.allow_multiple   = allow_multiple
        self.count_extensions = count_extensions or LINK_EXTENSIONS
        self.selected_paths   = []
        self.on_change        = None

        self.configure(
            highlightbackground=C["border"],
            highlightthickness=2,
            highlightcolor=C["primary"]
        )

        self.title_lbl = tk.Label(
            self, text="▸ " + label_text,
            font=("Helvetica", 12, "bold"),
            bg=C["surface"], fg=C["text"]
        )
        self.title_lbl.pack(pady=(12, 2), padx=14, anchor="w")

        self.hint_lbl = tk.Label(
            self, text=hint_text,
            font=("Helvetica", 9),
            bg=C["surface"], fg=C["sub"]
        )
        self.hint_lbl.pack(padx=14, anchor="w")

        pf = tk.Frame(self, bg=C["surface"])
        pf.pack(fill=tk.X, padx=14, pady=(6, 2))

        self.path_var = tk.StringVar()
        self.path_entry = tk.Entry(
            pf, textvariable=self.path_var,
            font=("Helvetica", 9),
            bg=C["input_bg"], fg=C["sub"],
            insertbackground=C["primary"],
            selectbackground=C["primary"],
            relief=tk.FLAT, bd=1
        )
        self.path_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
        self.path_entry.insert(0, "ここにパスをペースト...")
        self.path_entry.bind("<FocusIn>",  self._entry_in)
        self.path_entry.bind("<FocusOut>", self._entry_out)
        self.path_entry.bind("<Return>",   self._entry_submit)

        FlatButton(
            pf, text="読込", command=self._entry_submit,
            font=("Helvetica", 9),
            bg=C["accent"], fg="#FFFFFF",
            padx=10, pady=3
        ).pack(side=tk.LEFT)

        self.info_lbl = tk.Label(
            self, text="",
            font=("Helvetica", 9),
            bg=C["surface"], fg=C["info"],
            justify=tk.LEFT, anchor="w", wraplength=440
        )
        self.info_lbl.pack(padx=14, anchor="w", pady=(2, 0))

        bf = tk.Frame(self, bg=C["surface"])
        bf.pack(padx=14, pady=(6, 12), anchor="w")

        label = "ファイル選択" if select_mode == "file" else "フォルダ選択"
        self.sel_btn = FlatButton(
            bf, text=label, command=self._on_click,
            font=("Helvetica", 9),
            bg=C["accent"], fg="#FFFFFF",
            padx=12, pady=4
        )
        self.sel_btn.pack(side=tk.LEFT, padx=(0, 6))
        self.clr_btn = None

        for w in [self, self.title_lbl, self.hint_lbl]:
            w.bind("<Button-1>", lambda e: self._on_click())

        if _DND_BACKEND == "tkdnd":
            try:
                self.drop_target_register(DND_FILES)
                self.dnd_bind("<<Drop>>",      self._on_drop)
                self.dnd_bind("<<DragEnter>>", self._on_enter)
                self.dnd_bind("<<DragLeave>>", self._on_leave)
            except Exception:
                pass

    def _on_enter(self, event):
        self.configure(highlightbackground=C["primary"], highlightthickness=3,
                       bg=C["drop_hi"])
        for w in [self.title_lbl, self.hint_lbl, self.info_lbl]:
            w.configure(bg=C["drop_hi"])
        return event.action

    def _on_leave(self, event):
        bg = C["surface"]
        self.configure(bg=bg, highlightthickness=2,
                       highlightbackground=C["ok"] if self.selected_paths else C["border"])
        for w in [self.title_lbl, self.hint_lbl, self.info_lbl]:
            w.configure(bg=bg)
        return event.action

    def _on_drop(self, event):
        self._on_leave(event)
        paths = self._parse_paths(event.data)
        if paths:
            self._set_paths(paths)
        return event.action

    @staticmethod
    def _parse_paths(raw):
        paths, i = [], 0
        while i < len(raw):
            if raw[i] == '{':
                end = raw.index('}', i)
                paths.append(raw[i+1:end]); i = end + 2
            elif raw[i] == ' ':
                i += 1
            else:
                end = raw.find(' ', i)
                if end == -1: end = len(raw)
                paths.append(raw[i:end]); i = end + 1
        return [p for p in paths if p.strip()]

    def _entry_in(self, e=None):
        if self.path_var.get() == "ここにパスをペースト...":
            self.path_entry.delete(0, tk.END)
            self.path_entry.configure(fg=C["text"])

    def _entry_out(self, e=None):
        if not self.path_var.get().strip():
            self.path_entry.configure(fg=C["sub"])
            self.path_entry.insert(0, "ここにパスをペースト...")

    def _entry_submit(self, e=None):
        raw = self.path_var.get().strip()
        if not raw or raw == "ここにパスをペースト...": return
        path = raw.strip("'\"").strip()
        if os.path.exists(path):
            self._set_paths([path])
            self.path_entry.delete(0, tk.END)
            self._entry_out()
        else:
            self.path_entry.configure(fg=C["err"])

    def _on_click(self):
        if self.select_mode == "folder":
            p = filedialog.askdirectory(title="フォルダを選択")
            if p: self._set_paths([p])
        elif self.allow_multiple:
            ps = filedialog.askopenfilenames(
                title="ファイルを選択",
                filetypes=self.file_types or [("All", "*.*")])
            if ps: self._set_paths(list(ps))
        else:
            p = filedialog.askopenfilename(
                title="ファイルを選択",
                filetypes=self.file_types or [("All", "*.*")])
            if p: self._set_paths([p])

    def _set_paths(self, paths):
        if self.select_mode == "folder" and not self.allow_multiple:
            self.selected_paths = [paths[0]]
        elif self.allow_multiple:
            ex = set(self.selected_paths)
            for p in paths:
                if p not in ex:
                    self.selected_paths.append(p); ex.add(p)
        else:
            self.selected_paths = [paths[0]]
        self._update_display()
        if self.on_change: self.on_change()

    def _update_display(self):
        if not self.selected_paths:
            self.info_lbl.config(text="")
            self.configure(highlightbackground=C["border"])
            if self.clr_btn: self.clr_btn.destroy(); self.clr_btn = None
            return
        self.configure(highlightbackground=C["ok"], highlightthickness=2)
        names = [Path(p).name for p in self.selected_paths]
        if len(names) == 1 and self.select_mode == "folder":
            folder = Path(self.selected_paths[0])
            files = [f for f in folder.rglob("*")
                     if f.is_file() and f.suffix.lower() in self.count_extensions]
            text = f"✓  {folder.name}/  （{len(files)} ファイル）"
        elif len(names) <= 4:
            text = "\n".join(f"✓  {n}" for n in names)
        else:
            text = "\n".join(f"✓  {n}" for n in names[:3])
            text += f"\n  … 他 {len(names)-3} 件"
        self.info_lbl.config(text=text, fg=C["ok"])
        if not self.clr_btn:
            self.clr_btn = FlatButton(
                self.sel_btn.master, text="クリア", command=self._clear,
                font=("Helvetica", 8),
                bg="#3A1010", fg=C["err"],
                padx=8, pady=3
            )
            self.clr_btn.pack(side=tk.LEFT)

    def _clear(self):
        self.selected_paths = []
        self.info_lbl.config(fg=C["info"])
        self._update_display()
        if self.on_change: self.on_change()


# ════════════════════════════════════════════════════════════════
#  リンク処理コア
# ════════════════════════════════════════════════════════════════

import re as _re
import unicodedata as _ud

def _nfc(s):
    return _ud.normalize("NFC", s)

def _strip_number(name):
    name = name.translate(str.maketrans('０１２３４５６７８９', '0123456789'))
    name = _re.sub(r'^[\d\s\.\-_、。．・()（）【】\[\]「」『』〔〕]+', '', name)
    return name.strip()

def get_file_map(link_dirs):
    """リンク資料フォルダを走査してキー→相対パスのマップを返す"""
    fm = {}
    for link_dir in link_dirs:
        base = Path(link_dir)
        for f in sorted(base.rglob("*")):
            if f.is_file() and f.suffix.lower() in LINK_EXTENSIONS:
                rel_str = str(f.relative_to(base.parent)).replace("\\", "/")
                stem_nfc = _nfc(f.stem)
                name_nfc = _nfc(f.name)
                fm[stem_nfc] = rel_str
                fm[name_nfc] = rel_str
                stripped = _strip_number(stem_nfc)
                if stripped:
                    fm[stripped] = rel_str
                    fm[stripped + f.suffix] = rel_str
    return fm

def _to_file_uri(path: Path) -> str:
    """絶対パスを file:// URI に変換。スペース・日本語等をURLエンコード。
    これにより Word のハイパーリンクをクリック一発で開けるようにする。"""
    abs_str = str(path.resolve())
    if platform.system() == "Windows":
        abs_str = abs_str.replace("\\", "/")
        encoded = urllib.parse.quote(abs_str, safe="/:")
        return "file:///" + encoded.lstrip("/")
    else:
        encoded = urllib.parse.quote(abs_str, safe="/:")
        return "file://" + encoded

def copy_link_trees(link_dirs, dst_parent):
    for link_dir in link_dirs:
        src = Path(link_dir)
        dst = Path(dst_parent) / src.name
        if dst.exists(): shutil.rmtree(str(dst))
        shutil.copytree(str(src), str(dst))

def iter_all_paragraphs(doc):
    def _tables(tables):
        for t in tables:
            for row in t.rows:
                for cell in row.cells:
                    yield from cell.paragraphs
                    yield from _tables(cell.tables)
    yield from doc.paragraphs
    yield from _tables(doc.tables)

def _make_run(text, rpr=None):
    r = OxmlElement("w:r")
    if rpr is not None: r.append(deepcopy(rpr))
    t = OxmlElement("w:t"); t.text = text
    if text != text.strip(): t.set(qn("xml:space"), "preserve")
    r.append(t); return r

def _make_hyperlink(rid, text, rpr=None):
    hl = OxmlElement("w:hyperlink")
    hl.set(qn("r:id"), rid); hl.set(qn("w:history"), "1")
    r = OxmlElement("w:r")
    rp = OxmlElement("w:rPr")
    rs = OxmlElement("w:rStyle"); rs.set(qn("w:val"), "Hyperlink"); rp.append(rs)
    co = OxmlElement("w:color");  co.set(qn("w:val"), "0563C1");    rp.append(co)
    u  = OxmlElement("w:u");      u.set(qn("w:val"), "single");     rp.append(u)
    skip = {"rStyle", "color", "u"}
    if rpr is not None:
        for c in rpr:
            loc = c.tag.split("}")[-1] if "}" in c.tag else c.tag
            if loc not in skip: rp.append(deepcopy(c))
    r.append(rp)
    t = OxmlElement("w:t"); t.text = text
    if text != text.strip(): t.set(qn("xml:space"), "preserve")
    r.append(t); hl.append(r); return hl

def _find_matches(full, fm):
    keys = sorted(fm.keys(), key=len, reverse=True)
    matches, used = [], set()
    for k in keys:
        pos = 0
        while True:
            idx = full.find(k, pos)
            if idx == -1: break
            end = idx + len(k)
            span = set(range(idx, end))
            if not span & used:
                matches.append((idx, end, fm[k])); used |= span
            pos = idx + 1
    return sorted(matches, key=lambda x: x[0])

def process_paragraph(para, fm, part):
    p = para._p
    runs, pos = [], 0
    for ch in list(p):
        if ch.tag == qn("w:r"):
            te = ch.find(qn("w:t"))
            tx = (te.text or "") if te is not None else ""
            runs.append(dict(elem=ch, text=tx, start=pos,
                             end=pos+len(tx), rpr=ch.find(qn("w:rPr"))))
            pos += len(tx)
    if not runs: return 0
    full = "".join(r["text"] for r in runs)
    ms = _find_matches(_nfc(full), fm)
    if not ms: return 0
    crpr = [None] * len(full)
    for rd in runs:
        for i in range(rd["start"], rd["end"]): crpr[i] = rd["rpr"]
    for rd in runs: p.remove(rd["elem"])
    ppr = p.find(qn("w:pPr"))
    ins = (list(p).index(ppr)+1) if ppr is not None else 0
    elems, prev = [], 0
    for (s, e, rp) in ms:
        if prev < s:
            elems.append(_make_run(full[prev:s], crpr[prev] if prev < len(crpr) else None))
        rid = part.relate_to(rp, HYPERLINK_TYPE, is_external=True)
        elems.append(_make_hyperlink(rid, full[s:e], crpr[s] if s < len(crpr) else None))
        prev = e
    if prev < len(full):
        elems.append(_make_run(full[prev:], crpr[prev] if prev < len(crpr) else None))
    for i, el in enumerate(elems): p.insert(ins+i, el)
    return len(ms)


# ════════════════════════════════════════════════════════════════
#  PDF変換コア ─ ネイティブ Office 優先 / LibreOffice フォールバック
# ════════════════════════════════════════════════════════════════
#
#  優先順位:
#    Windows → Word/Excel/PowerPoint COM オートメーション（pywin32）
#    Mac     → AppleScript 経由でネイティブ Office
#    共通    → LibreOffice（未インストール時・非対応形式）
#    画像    → Pillow → LibreOffice の順でフォールバック
# ════════════════════════════════════════════════════════════════

# Mac で LibreOffice 変換時のフォント置換マップ
# Windows 専用日本語フォント → 游フォント（macOS 10.11+ 標準搭載、同系統で行間メトリクス近似）
# ※ HiraginoはMS明朝より行間メトリクスが大きく、w:line="0"と組み合わさると行間崩れの原因になる
_MAC_FONT_MAP = {
    # 明朝系 → 游明朝（MS明朝と同じFontworks設計、macOS標準搭載）
    "ＭＳ 明朝":      "游明朝",
    "ＭＳ Ｐ明朝":    "游明朝",
    "MS Mincho":      "游明朝",
    "MS PMincho":     "游明朝",
    # ゴシック系 → 游ゴシック（同系統）
    "ＭＳ ゴシック":   "游ゴシック",
    "ＭＳ Ｐゴシック": "游ゴシック",
    "MS Gothic":      "游ゴシック",
    "MS PGothic":     "游ゴシック",
    "メイリオ":        "游ゴシック",
    "Meiryo":         "游ゴシック",
}

import re as _re


def _patch_ms_fonts(src: Path) -> Path:
    r"""Mac 向け: Windowsフォント名を Mac 互換フォント名に置換した一時ファイルを返す。
    また w:line="0" w:lineRule="atLeast" を単一行間 (line=240, auto) に補正する。
    Word は line=0 を「自動行間」として扱うが LibreOffice は「0pt以上」と解釈し、
    フォントメトリクスが大きいと巨大な行間になるため修正が必要。"""
    import zipfile
    import tempfile
    tmp = Path(tempfile.mktemp(suffix=src.suffix))
    with zipfile.ZipFile(src, "r") as zin, \
         zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") or item.filename.endswith(".rels"):
                text = data.decode("utf-8", errors="replace")
                # ① フォント名置換
                for win_font, mac_font in _MAC_FONT_MAP.items():
                    text = text.replace(win_font, mac_font)
                # ② w:line="0" w:lineRule="atLeast" → 単一行間に補正
                text = _re.sub(
                    r'(w:line="0")(\s+)(w:lineRule="atLeast")',
                    r'w:line="240"\2w:lineRule="auto"', text)
                text = _re.sub(
                    r'(w:lineRule="atLeast")(\s+)(w:line="0")',
                    r'w:lineRule="auto"\2w:line="240"', text)
                data = text.encode("utf-8")
            zout.writestr(item, data)
    return tmp


# LibreOffice フォールバック用フィルタマップ
_LO_FILTER_MAP = {
    **{e: "writer_pdf_Export"  for e in _WORD_EXTS},
    **{e: "calc_pdf_Export"    for e in _EXCEL_EXTS},
    **{e: "impress_pdf_Export" for e in _PPT_EXTS},
    **{e: "draw_pdf_Export"    for e in _IMG_EXTS},
}


def _get_libreoffice_path():
    if platform.system() == "Windows":
        for p in [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]:
            if os.path.exists(p): return f'"{p}"'
        return "soffice"
    if platform.system() == "Darwin":
        mac = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        return mac if os.path.exists(mac) else "soffice"
    return "soffice"


def _convert_libreoffice(input_path: Path, output_dir: Path, ext: str):
    """LibreOffice によるフォールバック変換。
    Mac では Windows 専用フォントを Hiragino 系に置換してから変換する。"""
    output_dir.mkdir(parents=True, exist_ok=True)

    # Mac + Office ファイル → フォント置換した一時ファイルで変換
    patched: Path | None = None
    convert_path = input_path
    if platform.system() == "Darwin" and ext in (
        _WORD_EXTS | _EXCEL_EXTS | _PPT_EXTS
    ):
        try:
            patched = _patch_ms_fonts(input_path)
            convert_path = patched
        except Exception:
            pass  # 置換失敗時は元ファイルで続行

    lo = _get_libreoffice_path()
    lo_filter = _LO_FILTER_MAP.get(ext, "writer_pdf_Export")
    cmd = (
        f'{lo} --headless --norestore --nofirststartwizard '
        f'--convert-to "pdf:{lo_filter}" '
        f'"{convert_path}" --outdir "{output_dir}"'
    )
    try:
        result = subprocess.run(cmd, shell=True,
                                stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                                timeout=120)
        # 出力ファイル名は元のstemで確定（patchedファイルはtempなので）
        out_pdf = output_dir / (input_path.stem + ".pdf")
        if not out_pdf.exists() and patched:
            # LibreOffice は変換元ファイル名で出力するため、temp名で生成された可能性
            tmp_pdf = output_dir / (patched.stem + ".pdf")
            if tmp_pdf.exists():
                tmp_pdf.rename(out_pdf)
        if out_pdf.exists():
            return True, str(out_pdf)
        err = result.stderr.decode(errors="ignore")
        return False, err or "出力ファイルが生成されませんでした"
    except subprocess.TimeoutExpired:
        return False, "タイムアウト（120秒）"
    except Exception as e:
        return False, str(e)
    finally:
        if patched and patched.exists():
            try:
                patched.unlink()
            except Exception:
                pass


def _convert_image(input_path: Path, output_dir: Path):
    """画像 → PDF。Pillow があれば使用、なければ LibreOffice へ"""
    try:
        from PIL import Image
        output_dir.mkdir(parents=True, exist_ok=True)
        out_pdf = output_dir / (input_path.stem + ".pdf")
        with Image.open(input_path) as img:
            img.convert("RGB").save(str(out_pdf), "PDF", resolution=150)
        if out_pdf.exists():
            return True, str(out_pdf)
        return False, "画像PDF変換失敗"
    except ImportError:
        return _convert_libreoffice(input_path, output_dir, input_path.suffix.lower())
    except Exception as e:
        return False, str(e)


# ── Windows COM オートメーション ──────────────────────────────

def _win_word(abs_in: str, abs_out: str):
    """Word COM: .doc/.docx/.odt/.rtf/.txt → PDF"""
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return None, "pywin32 未インストール"
    pythoncom.CoInitialize()
    word = None; doc = None
    try:
        word = win32com.client.DispatchEx("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0          # wdAlertsNone
        doc = word.Documents.Open(
            abs_in,
            ConfirmConversions=False,
            ReadOnly=True,
            AddToRecentFiles=False,
            NoEncodingDialog=True,
        )
        doc.ExportAsFixedFormat(
            OutputFileName=abs_out,
            ExportFormat=17,            # wdExportFormatPDF
            OpenAfterExport=False,
            OptimizeFor=0,              # wdExportOptimizeForPrint
            Range=0,                    # wdExportAllDocument
            IncludeDocProps=True,
            KeepIRM=True,
            CreateBookmarks=0,
            DocStructureTags=True,
            BitmapMissingFonts=True,
            UseISO19005_1=False,
        )
        out = Path(abs_out)
        return (True, str(out)) if out.exists() else (False, "PDF未生成")
    except Exception as e:
        return False, str(e)
    finally:
        if doc:
            try: doc.Close(False)
            except: pass
        if word:
            try: word.Quit()
            except: pass
        pythoncom.CoUninitialize()


def _win_excel(abs_in: str, abs_out: str):
    """Excel COM: .xls/.xlsx/.ods/.csv → PDF"""
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return None, "pywin32 未インストール"
    pythoncom.CoInitialize()
    excel = None; wb = None
    try:
        excel = win32com.client.DispatchEx("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        excel.AskToUpdateLinks = False
        wb = excel.Workbooks.Open(
            abs_in, UpdateLinks=0, ReadOnly=True, AddToMru=False,
        )
        wb.ExportAsFixedFormat(
            Type=0,                     # xlTypePDF
            Filename=abs_out,
            Quality=0,                  # xlQualityStandard
            IncludeDocProperties=True,
            IgnorePrintAreas=False,
            OpenAfterPublish=False,
        )
        out = Path(abs_out)
        return (True, str(out)) if out.exists() else (False, "PDF未生成")
    except Exception as e:
        return False, str(e)
    finally:
        if wb:
            try: wb.Close(False)
            except: pass
        if excel:
            try: excel.Quit()
            except: pass
        pythoncom.CoUninitialize()


def _win_ppt(abs_in: str, abs_out: str):
    """PowerPoint COM: .ppt/.pptx/.odp → PDF"""
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return None, "pywin32 未インストール"
    pythoncom.CoInitialize()
    ppt = None; prs = None
    try:
        ppt = win32com.client.DispatchEx("PowerPoint.Application")
        prs = ppt.Presentations.Open(
            abs_in, ReadOnly=True, Untitled=True, WithWindow=False,
        )
        prs.ExportAsFixedFormat(
            Path=abs_out,
            FixedFormatType=2,          # ppFixedFormatTypePDF
            Intent=1,                   # ppFixedFormatIntentPrint
            HandoutOrder=1,
            OutputType=1,               # ppPrintOutputSlides
            PrintHiddenSlides=False,
            PrintRange=None,
            RangeType=1,                # ppPrintAll
            SlideShowName="",
            IncludeDocProperties=True,
            KeepIRM=True,
            DocStructureTags=True,
            BitmapMissingFonts=True,
            UseISO19005_1=False,
        )
        out = Path(abs_out)
        return (True, str(out)) if out.exists() else (False, "PDF未生成")
    except Exception as e:
        return False, str(e)
    finally:
        if prs:
            try: prs.Close()
            except: pass
        if ppt:
            try: ppt.Quit()
            except: pass
        pythoncom.CoUninitialize()


def _convert_windows(input_path: Path, output_dir: Path, ext: str):
    """Windows: COM オートメーションでネイティブ変換、失敗時は LibreOffice"""
    output_dir.mkdir(parents=True, exist_ok=True)
    abs_in  = str(input_path.resolve())
    abs_out = str((output_dir / (input_path.stem + ".pdf")).resolve())

    if ext in _WORD_EXTS:
        result = _win_word(abs_in, abs_out)
    elif ext in _EXCEL_EXTS:
        result = _win_excel(abs_in, abs_out)
    elif ext in _PPT_EXTS:
        result = _win_ppt(abs_in, abs_out)
    elif ext in _IMG_EXTS:
        return _convert_image(input_path, output_dir)
    else:
        return _convert_libreoffice(input_path, output_dir, ext)

    # result[0] が None = pywin32 なし → LibreOffice にフォールバック
    if result[0] is None:
        return _convert_libreoffice(input_path, output_dir, ext)
    return result


# ── Mac AppleScript ───────────────────────────────────────────

def _esc_as(s: str) -> str:
    r"""AppleScript 文字列リテラル用エスケープ（" と \ を処理）"""
    return s.replace("\\", "\\\\").replace('"', '\\"')


def _run_applescript(script: str, timeout: int = 120):
    """osascript でスクリプトを stdin から実行。(returncode, stderr) を返す"""
    try:
        result = subprocess.run(
            ["osascript"],
            input=script.encode("utf-8"),
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            timeout=timeout,
        )
        return result.returncode, result.stderr.decode(errors="ignore")
    except subprocess.TimeoutExpired:
        return -1, "タイムアウト（120秒）"
    except Exception as e:
        return -1, str(e)


def _mac_word(abs_in: str, abs_out: str):
    """Mac: AppleScript で Word → PDF"""
    if not Path("/Applications/Microsoft Word.app").exists():
        return None, "Microsoft Word 未インストール"
    script = f'''
tell application "Microsoft Word"
    try
        open POSIX file "{_esc_as(abs_in)}"
        delay 2
        try
            save as (active document) file name "{_esc_as(abs_out)}" file format format PDF
        end try
        close (active document) saving no
    on error errMsg
        try
            close (active document) saving no
        end try
    end try
end tell
'''
    rc, err = _run_applescript(script)
    out = Path(abs_out)
    if out.exists():
        return True, str(out)
    return False, err or "PDF未生成（Word AppleScript）"


def _mac_excel(abs_in: str, abs_out: str):
    """Mac: AppleScript で Excel → PDF"""
    if not Path("/Applications/Microsoft Excel.app").exists():
        return None, "Microsoft Excel 未インストール"
    script = f'''
tell application "Microsoft Excel"
    try
        open POSIX file "{_esc_as(abs_in)}"
        delay 2
        try
            save workbook as (active workbook) filename "{_esc_as(abs_out)}" file format PDF file format
        end try
        close (active workbook) saving no
    on error errMsg
        try
            close (active workbook) saving no
        end try
    end try
end tell
'''
    rc, err = _run_applescript(script)
    out = Path(abs_out)
    if out.exists():
        return True, str(out)
    return False, err or "PDF未生成（Excel AppleScript）"


def _mac_ppt(abs_in: str, abs_out: str):
    """Mac: AppleScript で PowerPoint → PDF"""
    if not Path("/Applications/Microsoft PowerPoint.app").exists():
        return None, "Microsoft PowerPoint 未インストール"
    script = f'''
tell application "Microsoft PowerPoint"
    try
        open POSIX file "{_esc_as(abs_in)}"
        delay 2
        try
            save active presentation in "{_esc_as(abs_out)}" as save as PDF
        end try
        close active presentation saving no
    on error errMsg
        try
            close active presentation saving no
        end try
    end try
end tell
'''
    rc, err = _run_applescript(script)
    out = Path(abs_out)
    if out.exists():
        return True, str(out)
    return False, err or "PDF未生成（PowerPoint AppleScript）"


def _convert_mac(input_path: Path, output_dir: Path, ext: str):
    """Mac: ネイティブ Office AppleScript、失敗時は LibreOffice"""
    output_dir.mkdir(parents=True, exist_ok=True)
    abs_in  = str(input_path.resolve())
    abs_out = str((output_dir / (input_path.stem + ".pdf")).resolve())

    if ext in _WORD_EXTS:
        result = _mac_word(abs_in, abs_out)
    elif ext in _EXCEL_EXTS:
        result = _mac_excel(abs_in, abs_out)
    elif ext in _PPT_EXTS:
        result = _mac_ppt(abs_in, abs_out)
    elif ext in _IMG_EXTS:
        return _convert_image(input_path, output_dir)
    else:
        return _convert_libreoffice(input_path, output_dir, ext)

    # result[0] が None = Office 未インストール
    # result[0] が False = AppleScript 失敗
    # → どちらも LibreOffice にフォールバック
    if not result[0]:
        return _convert_libreoffice(input_path, output_dir, ext)
    return result


def convert_to_pdf(input_path: Path, output_dir: Path, log_cb=None):
    """メイン変換エントリ。OS に応じてネイティブ Office を優先し、
    未インストール・失敗時は LibreOffice へ自動フォールバック。"""
    output_dir.mkdir(parents=True, exist_ok=True)
    ext = input_path.suffix.lower()
    os_name = platform.system()

    if os_name == "Windows":
        return _convert_windows(input_path, output_dir, ext)
    elif os_name == "Darwin":
        return _convert_mac(input_path, output_dir, ext)
    else:
        return _convert_libreoffice(input_path, output_dir, ext)


def scan_all_pdf_targets(paths):
    """PDF変換・コピー対象を再帰収集。
    PDF_EXTENSIONS → 変換対象、.pdf → コピー対象として収集する。"""
    result = []
    for p in paths:
        p = Path(p)
        if p.is_dir():
            for f in sorted(p.rglob("*")):
                if f.is_file() and (f.suffix.lower() in PDF_EXTENSIONS
                                    or f.suffix.lower() == ".pdf"):
                    result.append(f)
        elif p.is_file():
            if p.suffix.lower() in PDF_EXTENSIONS or p.suffix.lower() == ".pdf":
                result.append(p)
    return result


# ════════════════════════════════════════════════════════════════
#  共通ヘッダーウィジェット
# ════════════════════════════════════════════════════════════════

def build_header(parent, title, subtitle, show_version=True, on_update=None):
    hdr = tk.Frame(parent, bg=C["accent"], height=56)
    hdr.pack(fill=tk.X); hdr.pack_propagate(False)
    tk.Label(hdr, text=title,
             font=("Helvetica", 18, "bold"),
             bg=C["accent"], fg=C["primary"]
             ).pack(side=tk.LEFT, padx=20, pady=10)
    tk.Label(hdr, text=subtitle,
             font=("Helvetica", 10),
             bg=C["accent"], fg="#C8DCF5"
             ).pack(side=tk.LEFT, pady=10)
    if show_version:
        if on_update:
            FlatButton(hdr, text="🔄", command=on_update,
                       font=("Helvetica", 11),
                       bg=C["accent"], fg="#AACFEE",
                       padx=6, pady=0
                       ).pack(side=tk.RIGHT, pady=10)
        tk.Label(hdr, text=f"v{APP_VERSION}",
                 font=("Helvetica", 8),
                 bg=C["accent"], fg="#AACFEE"
                 ).pack(side=tk.RIGHT, padx=(0, 2), pady=10)
    tk.Frame(parent, bg=C["primary"], height=3).pack(fill=tk.X)
    return hdr


# ════════════════════════════════════════════════════════════════
#  ランチャー画面
# ════════════════════════════════════════════════════════════════

class LauncherFrame(tk.Frame):

    def __init__(self, parent, on_link, on_pdf, on_checker):
        super().__init__(parent, bg=C["bg"])
        self.on_link    = on_link
        self.on_pdf     = on_pdf
        self.on_checker = on_checker
        self._build()

    def _build(self):
        build_header(self, "⛓  楽々JC", "自動処理ツール集")
        center = tk.Frame(self, bg=C["bg"])
        center.pack(expand=True)
        tk.Label(center, text="使用する機能を選んでください",
                 font=("Helvetica", 13),
                 bg=C["bg"], fg=C["text"]
                 ).pack(pady=(40, 32))
        btn_area = tk.Frame(center, bg=C["bg"])
        btn_area.pack()
        self._feature_btn(btn_area, "⚓", "統一ルール修正",
                          "Word / Excel / PPT / PDF\n統一語句チェック・修正",
                          self.on_checker, col=0)
        self._feature_btn(btn_area, "📄", "PDF一括変換",
                          "Word / Excel / PowerPoint等を\nPDFに一括変換",
                          self.on_pdf, col=1)
        self._feature_btn(btn_area, "⛓", "リンク一括設定",
                          "Wordファイルへ\nハイパーリンクを自動挿入",
                          self.on_link, col=2)
        tk.Label(center, text=f"楽々JC  v{APP_VERSION}",
                 font=("Helvetica", 8),
                 bg=C["bg"], fg="#5A7AAA"
                 ).pack(pady=(40, 0))

    def _feature_btn(self, parent, icon, title, desc, command, col):
        frame = tk.Frame(parent, bg=C["surface"],
                         highlightbackground=C["border"],
                         highlightthickness=2, cursor="hand2")
        frame.grid(row=0, column=col, padx=14, pady=4, ipadx=6, ipady=6)
        tk.Label(frame, text=icon, font=("Helvetica", 36),
                 bg=C["surface"], fg=C["primary"]).pack(pady=(24, 6), padx=36)
        tk.Label(frame, text=title, font=("Helvetica", 14, "bold"),
                 bg=C["surface"], fg=C["text"]).pack()
        tk.Label(frame, text=desc, font=("Helvetica", 9),
                 bg=C["surface"], fg=C["sub"],
                 justify=tk.CENTER).pack(pady=(4, 20), padx=20)
        def on_enter(e, f=frame):
            f.configure(highlightbackground=C["primary"])
        def on_leave(e, f=frame):
            f.configure(highlightbackground=C["border"])
        def on_click(e): command()
        for w in frame.winfo_children() + [frame]:
            w.bind("<Enter>", on_enter)
            w.bind("<Leave>", on_leave)
            w.bind("<Button-1>", on_click)

    def _placeholder_btn(self, parent, col):
        frame = tk.Frame(parent, bg=C["surface"],
                         highlightbackground=C["border"], highlightthickness=1)
        frame.grid(row=0, column=col, padx=14, pady=4, ipadx=6, ipady=6)
        tk.Label(frame, text="＋", font=("Helvetica", 36),
                 bg=C["surface"], fg=C["border"]).pack(pady=(24, 6), padx=36)
        tk.Label(frame, text="機能追加予定", font=("Helvetica", 12),
                 bg=C["surface"], fg=C["border"]).pack()
        tk.Label(frame, text="近日公開", font=("Helvetica", 9),
                 bg=C["surface"], fg=C["border"]).pack(pady=(4, 20), padx=20)


# ════════════════════════════════════════════════════════════════
#  リンク一括設定 画面
# ════════════════════════════════════════════════════════════════

class LinkFrame(LoggedFrame):

    def __init__(self, parent, on_back, on_go_pdf, on_go_checker, dnd_ok,
                 on_check_update=None):
        super().__init__(parent, bg=C["bg"])
        self.on_back         = on_back
        self.on_go_pdf       = on_go_pdf
        self.on_go_checker   = on_go_checker
        self.dnd_ok          = dnd_ok
        self.on_check_update = on_check_update
        self.folder_name_entries = []
        self._build()

    def _build(self):
        build_header(self, "⛓  リンク一括設定", "Word ハイパーリンク自動挿入",
                     on_update=self.on_check_update)
        nav = tk.Frame(self, bg=C["bg"])
        nav.pack(fill=tk.X, padx=16, pady=(8, 0))
        nav_button(nav, "← ホームへ戻る", self.on_back).pack(side=tk.LEFT)
        nav_button(nav, "PDF一括変換へ →",     self.on_go_pdf).pack(side=tk.RIGHT)
        nav_button(nav, "統一語句チェックへ →", self.on_go_checker).pack(
            side=tk.RIGHT, padx=(0, 6))
        self.dnd_lbl = tk.Label(nav,
                                text="D&D ✓" if self.dnd_ok else "---",
                                font=("Helvetica", 9), bg=C["bg"],
                                fg=C["ok"] if self.dnd_ok else "#AACFEE")
        self.dnd_lbl.pack(side=tk.RIGHT, padx=10)

        _, main = make_scrollable_frame(self)
        pad = dict(padx=16, pady=(0, 10))

        self.word_zone = DropZone(
            main, "計画書（Word ファイル）",
            "ドラッグ＆ドロップ、またはクリックして選択（複数可）",
            select_mode="file",
            file_types=[("Word 文書", "*.docx")],
            allow_multiple=True
        )
        self.word_zone.pack(fill=tk.X, **pad)
        self.word_zone.on_change = self._on_word_changed

        self.link_zone = DropZone(
            main, "リンク資料フォルダ",
            "ドラッグ＆ドロップで複数追加可能（何階層でも対応）",
            select_mode="folder", allow_multiple=True
        )
        self.link_zone.pack(fill=tk.X, **pad)
        self.link_zone.on_change = self._check_ready

        self.output_zone = DropZone(
            main, "出力先フォルダ",
            "ドラッグ＆ドロップ、またはクリックして選択",
            select_mode="folder", allow_multiple=False
        )
        self.output_zone.pack(fill=tk.X, **pad)
        self.output_zone.on_change = self._check_ready

        self.fname_outer = tk.Frame(main, bg=C["bg"])
        self.structure_lbl = tk.Label(
            main, text="", font=("Helvetica", 9),
            bg=C["bg"], fg=C["sub"], justify=tk.LEFT, anchor="w"
        )
        self.structure_lbl.pack(padx=16, anchor="w")

        section_divider(main)

        self.run_btn = FlatButton(
            main, text="▶  リンクを作成", command=self._run,
            font=("Helvetica", 14), bold=True,
            bg=C["accent"], fg="#AACFEE",
            padx=30, pady=12, cursor="arrow", state=tk.DISABLED
        )
        self.run_btn.pack(pady=(0, 10))

        self.log_txt = make_log_widget(main)
        log_write(self.log_txt,
                  "D&D: 有効 ✓" if self.dnd_ok else "D&D: ダイアログ / パスペーストを利用",
                  "success" if self.dnd_ok else "info")

    def _on_word_changed(self):
        self._rebuild_folder_names()
        self._check_ready()

    def _rebuild_folder_names(self):
        self.fname_outer.pack_forget()
        for w in self.fname_outer.winfo_children(): w.destroy()
        self.folder_name_entries = []
        word_paths = self.word_zone.selected_paths
        if not word_paths:
            self.structure_lbl.configure(text=""); return
        tk.Label(self.fname_outer,
                 text="出力フォルダ名（変更可）",
                 font=("Helvetica", 10, "bold"),
                 bg=C["bg"], fg=C["sub"]
                 ).pack(anchor="w", padx=16, pady=(0, 4))
        for wp in word_paths:
            p = Path(wp)
            if p.name.startswith("~$"): continue
            row = tk.Frame(self.fname_outer, bg=C["bg"])
            row.pack(fill=tk.X, padx=16, pady=2)
            tk.Label(row, text=f"{p.name}  →",
                     font=("Helvetica", 9), bg=C["bg"], fg=C["sub"],
                     anchor="e", width=30).pack(side=tk.LEFT, padx=(0, 6))
            var = tk.StringVar(value=p.stem)
            tk.Entry(row, textvariable=var,
                     font=("Helvetica", 10),
                     bg=C["input_bg"], fg=C["text"],
                     insertbackground=C["primary"],
                     relief=tk.FLAT, bd=1
                     ).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 6))
            def _reset(v=var, d=p.stem): v.set(d)
            FlatButton(row, text="戻す", command=_reset,
                       font=("Helvetica", 8),
                       bg=C["accent"], fg="#FFFFFF",
                       padx=6, pady=2).pack(side=tk.LEFT)
            self.folder_name_entries.append((wp, var))
        self.fname_outer.pack(fill=tk.X, pady=(0, 4), before=self.structure_lbl)
        self._update_structure_label()

    def _update_structure_label(self):
        if not self.folder_name_entries: return
        lines = ["📁 出力フォルダ構成プレビュー:"]
        for wp, var in self.folder_name_entries:
            fn = var.get() or Path(wp).stem
            lines.append(f"  出力先/{fn}/")
            lines.append(f"    ├─ {Path(wp).name}")
            for lp in self.link_zone.selected_paths:
                lines.append(f"    ├─ {Path(lp).name}/")
        self.structure_lbl.configure(text="\n".join(lines))

    def _check_ready(self):
        self._update_structure_label()
        ok = (self.word_zone.selected_paths and
              self.link_zone.selected_paths and
              self.output_zone.selected_paths)
        if ok:
            self.run_btn.configure(state=tk.NORMAL, bg=C["primary"],
                                   fg="white", cursor="hand2")
        else:
            self.run_btn.configure(state=tk.DISABLED, bg=C["accent"],
                                   fg="#AACFEE", cursor="arrow")

    def _run(self):
        self.run_btn.configure(state=tk.DISABLED,
                               text="⏳  処理中...", bg=C["accent"], fg=C["warn"])
        self.log_txt.configure(state=tk.NORMAL)
        self.log_txt.delete("1.0", tk.END)
        self.log_txt.configure(state=tk.DISABLED)
        threading.Thread(target=self._worker, daemon=True).start()

    def _worker(self):
        try:
            link_dirs  = self.link_zone.selected_paths
            output_dir = Path(self.output_zone.selected_paths[0])
            fn_map     = {wp: var.get().strip() or Path(wp).stem
                          for wp, var in self.folder_name_entries}

            fm = get_file_map(link_dirs)
            if not fm:
                self._log("[警告] リンク資料フォルダに対象ファイルがありません。", "warning")
                self.after(0, self._reset_btn); return

            self._log("リンク対象ファイル:", "info")
            seen = set()
            for v in fm.values():
                if v not in seen:
                    self._log(f"  {v}"); seen.add(v)
            self._log("")

            entries = [(wp, fn) for wp, fn in fn_map.items()
                       if not Path(wp).name.startswith("~$")]
            self._log(f"{len(entries)} 件の Word を処理...", "info")
            total = 0

            for wp, cname in entries:
                wpath = Path(wp)
                self._log(f"  {wpath.name}  →  {cname}/")
                try:
                    out = output_dir / cname
                    out.mkdir(parents=True, exist_ok=True)

                    # ① 先にリンクフォルダをコピー
                    copy_link_trees(link_dirs, out)
                    self._log(f"    資料フォルダ {len(link_dirs)} 個をコピー完了")

                    # ② 絶対 file:// URI マップを構築（クリック即開き保証）
                    #    fm の値は "FolderName/file.pdf" 形式 (out 直下からの相対)
                    abs_fm = {key: _to_file_uri(out / rel)
                              for key, rel in fm.items()}

                    # ③ Wordを開いてリンク挿入 → 保存
                    doc = Document(wpath)
                    cnt = sum(process_paragraph(p, abs_fm, doc.part)
                              for p in iter_all_paragraphs(doc))
                    doc.save(out / wpath.name)

                    if cnt > 0:
                        self._log(f"    {cnt} 箇所リンク挿入", "success")
                    else:
                        self._log("    対象テキストなし", "warning")
                    total += cnt
                except Exception as e:
                    self._log(f"    エラー: {e}", "error")

            self._log("")
            self._log(f"完了！  合計 {total} 箇所のリンクを挿入しました。", "success")
            self._log(f"出力先: {output_dir}", "info")
            self.after(0, lambda: messagebox.showinfo(
                "完了",
                f"処理が完了しました！\n\n合計 {total} 箇所のリンクを挿入\n出力先: {output_dir}"
            ))
        except Exception as e:
            self._log(f"\n予期しないエラー: {e}", "error")
        finally:
            self.after(0, self._reset_btn)

    def _reset_btn(self):
        self.run_btn.configure(text="▶  リンクを作成")
        self._check_ready()


# ════════════════════════════════════════════════════════════════
#  PDF一括変換 画面
# ════════════════════════════════════════════════════════════════

class PdfFrame(LoggedFrame):

    def __init__(self, parent, on_back, on_go_link, on_go_checker):
        super().__init__(parent, bg=C["bg"])
        self.on_back       = on_back
        self.on_go_link    = on_go_link
        self.on_go_checker = on_go_checker
        self._cancel_flag = threading.Event()
        self._build()

    def _build(self):
        build_header(self, "📄  PDF一括変換",
                     "Word / Excel / PowerPoint / 画像 → PDF")
        nav = tk.Frame(self, bg=C["bg"])
        nav.pack(fill=tk.X, padx=16, pady=(8, 0))
        nav_button(nav, "← ホームへ戻る", self.on_back).pack(side=tk.LEFT)
        nav_button(nav, "リンク一括設定へ →",    self.on_go_link).pack(side=tk.RIGHT)
        nav_button(nav, "統一語句チェックへ →", self.on_go_checker).pack(
            side=tk.RIGHT, padx=(0, 6))

        _, main = make_scrollable_frame(self)
        pad = dict(padx=16, pady=(0, 10))

        self.src_zone = DropZone(
            main, "変換対象フォルダ / ファイル",
            "複数フォルダをドラッグ&ドロップ（サブフォルダも自動検索）",
            select_mode="folder", allow_multiple=True,
            count_extensions=PDF_EXTENSIONS
        )
        self.src_zone.pack(fill=tk.X, **pad)
        self.src_zone.on_change = self._update_count

        # ── 出力先（常時表示・常時D&D対応）──
        self.custom_out_zone = DropZone(
            main, "出力先",
            "省略時は元ファイルと同じフォルダに保存  /  フォルダをドロップして出力先を指定",
            select_mode="folder", allow_multiple=False
        )
        self.custom_out_zone.pack(fill=tk.X, **pad)

        # ── 変換エンジン表示 ──
        _os = platform.system()
        if _os == "Darwin":
            _engine_hint = "変換: ネイティブ Office（AppleScript）優先 / LibreOffice フォールバック"
        elif _os == "Windows":
            _engine_hint = "変換: ネイティブ Office（COM）優先 / LibreOffice フォールバック"
        else:
            _engine_hint = "変換: LibreOffice"
        tk.Label(main, text=_engine_hint,
                 font=("Helvetica", 8), bg=C["bg"], fg="#5A7AAA"
                 ).pack(padx=16, anchor="w")

        self.count_lbl = tk.Label(main, text="",
                                  font=("Helvetica", 9), bg=C["bg"], fg=C["sub"])
        self.count_lbl.pack(padx=16, anchor="w")

        section_divider(main)

        btn_row = tk.Frame(main, bg=C["bg"])
        btn_row.pack(pady=(0, 10))

        self.run_btn = FlatButton(
            btn_row, text="▶  PDF変換を開始", command=self._run,
            font=("Helvetica", 14), bold=True,
            bg=C["accent"], fg="#AACFEE",
            padx=30, pady=12, cursor="arrow", state=tk.DISABLED
        )
        self.run_btn.pack(side=tk.LEFT, padx=(0, 12))

        self.cancel_btn = FlatButton(
            btn_row, text="■ 中断", command=self._cancel,
            font=("Helvetica", 11),
            bg="#3A1010", fg=C["err"],
            padx=16, pady=12, cursor="arrow", state=tk.DISABLED
        )
        self.cancel_btn.pack(side=tk.LEFT)

        self.log_txt = make_log_widget(main)
        log_write(self.log_txt, "フォルダをドロップして変換を開始してください", "info")
        log_write(self.log_txt, _engine_hint, "info")

        lo = _get_libreoffice_path()
        lo_exists = os.path.exists(lo.strip('"')) if lo != "soffice" else True
        if not lo_exists:
            log_write(self.log_txt,
                      "⚠️  LibreOffice が見つかりません（フォールバック不可）",
                      "warning")

    def _update_count(self):
        paths = self.src_zone.selected_paths
        if not paths:
            self.count_lbl.configure(text="")
            self.run_btn.configure(state=tk.DISABLED, bg=C["accent"],
                                   fg="#AACFEE", cursor="arrow")
            return
        files = scan_all_pdf_targets(paths)
        conv  = sum(1 for f in files if f.suffix.lower() in PDF_EXTENSIONS)
        copy_ = sum(1 for f in files if f.suffix.lower() == ".pdf")
        text  = f"対象ファイル: {len(files)} 件"
        if copy_:
            text += f"  （変換 {conv}、コピー {copy_}）"
        self.count_lbl.configure(text=text,
                                  fg=C["info"] if files else C["warn"])
        if files:
            self.run_btn.configure(state=tk.NORMAL, bg=C["primary"],
                                   fg="white", cursor="hand2")
        else:
            self.run_btn.configure(state=tk.DISABLED, bg=C["accent"],
                                   fg="#AACFEE", cursor="arrow")

    def _cancel(self):
        self._cancel_flag.set()
        self._log("中断リクエストを送信しました...", "warning")

    def _run(self):
        self._cancel_flag.clear()
        self.run_btn.configure(state=tk.DISABLED,
                               text="⏳  変換中...", bg=C["accent"], fg=C["warn"])
        self.cancel_btn.configure(state=tk.NORMAL, cursor="hand2")
        self.log_txt.configure(state=tk.NORMAL)
        self.log_txt.delete("1.0", tk.END)
        self.log_txt.configure(state=tk.DISABLED)
        threading.Thread(target=self._worker, daemon=True).start()

    def _worker(self):
        try:
            src_paths   = self.src_zone.selected_paths
            base_custom = (Path(self.custom_out_zone.selected_paths[0])
                           if self.custom_out_zone.selected_paths else None)

            all_files = []  # [(file_path, src_root)]
            for sp in src_paths:
                sp = Path(sp)
                if sp.is_dir():
                    for f in sorted(sp.rglob("*")):
                        if f.is_file():
                            ext = f.suffix.lower()
                            if ext in PDF_EXTENSIONS or ext == ".pdf":
                                all_files.append((f, sp))
                elif sp.is_file():
                    ext = sp.suffix.lower()
                    if ext in PDF_EXTENSIONS or ext == ".pdf":
                        all_files.append((sp, sp.parent))

            total = len(all_files)
            self._log(f"対象ファイル {total} 件を処理します...", "info")

            success, fail, skipped = 0, 0, 0
            for i, (f, src_root) in enumerate(all_files, 1):
                if self._cancel_flag.is_set():
                    self._log(f"中断しました（{i-1}/{total} 件処理済）", "warning")
                    break

                ext = f.suffix.lower()
                try:
                    rel = f.relative_to(src_root)
                except ValueError:
                    rel = Path(f.name)

                if base_custom:
                    # src_root.name を挟むことで、複数フォルダを同時投入しても
                    # 出力先に「ドロップ元フォルダ名/階層」が再現される
                    out_dir = base_custom / src_root.name / rel.parent
                else:
                    out_dir = f.parent

                out_dir.mkdir(parents=True, exist_ok=True)

                if ext == ".pdf":
                    dst = out_dir / f.name
                    try:
                        if dst.resolve() == f.resolve():
                            self._log(f"[{i}/{total}] スキップ（同一パス）: {f.name}", "info")
                            skipped += 1
                            continue
                    except Exception:
                        pass
                    self._log(f"[{i}/{total}] コピー: {f.name}")
                    try:
                        shutil.copy2(f, dst)
                        success += 1
                        self._log(f"  ✓ → {dst.name}", "success")
                    except Exception as e:
                        fail += 1
                        self._log(f"  ✕ {e}", "error")
                else:
                    out_pdf = out_dir / (f.stem + ".pdf")
                    if out_pdf.exists():
                        self._log(f"[{i}/{total}] スキップ（既存）: {f.stem}.pdf", "info")
                        skipped += 1
                        continue
                    self._log(f"[{i}/{total}] 変換: {f.name}")
                    ok, detail = convert_to_pdf(f, out_dir)
                    if ok:
                        success += 1
                        self._log(f"  ✓ → {f.stem}.pdf", "success")
                    else:
                        fail += 1
                        self._log(f"  ✕ {detail}", "error")

            self._log("")
            self._log(
                f"完了！  成功 {success} 件 ／ スキップ {skipped} 件 ／ 失敗 {fail} 件",
                "success"
            )
            self.after(0, lambda: messagebox.showinfo(
                "PDF変換完了",
                f"処理が完了しました！\n\n"
                f"成功: {success} 件\nスキップ: {skipped} 件\n失敗: {fail} 件"
            ))
        except Exception as e:
            self._log(f"予期しないエラー: {e}", "error")
        finally:
            self.after(0, self._reset_btns)

    def _reset_btns(self):
        self.cancel_btn.configure(state=tk.DISABLED, cursor="arrow")
        self.run_btn.configure(text="▶  PDF変換を開始")
        self._update_count()


# ════════════════════════════════════════════════════════════════
#  統一ルール修正チェッカー  ─  コアロジック
# ════════════════════════════════════════════════════════════════

def _checker_load_rules(csv_path):
    """rules.csv を読み込み {類義語: 統一語句} の辞書を返す（長さ降順）"""
    result = {}
    for enc in ('utf-8-sig', 'shift-jis', 'utf-8', 'cp932'):
        try:
            with open(csv_path, encoding=enc, newline='') as f:
                reader = csv.DictReader(f)
                rows = list(reader)
                fields = reader.fieldnames or []
            if '類義語' in fields and '統一語句' in fields:
                rows = sorted(rows, key=lambda r: len(r.get('類義語', '')), reverse=True)
                for row in rows:
                    k = row['類義語'].strip()
                    v = row['統一語句'].strip()
                    if k and v:
                        result[k] = v
                return result
        except Exception:
            continue
    return result


_CHECKER_KEEP_WORDS = [
    "会員に成長する機会", "会員拡大運動", "会員拡大", "正会員", "新入会員",
    "日本の青年会議所は", "希望をもたらす変革の起点として",
    "輝く個性が調和する未来を描き", "社会の課題を解決することで",
    "持続可能な地域を創ることを誓う", "われわれ JAYCEE は", "われわれJAYCEEは",
    "志高き組織ビジョン", "志高き人材育成ビジョン", "志高きまち創造ビジョン",
]

_ZEN_ALNUM = "ＡＢＣＤＥＦＧＨＩＪＫＬＭＮＯＰＱＲＳＴＵＶＷＸＹＺａｂｃｄｅｆｇｈｉｊｋｌｍｎｏｐｑｒｓｔｕｖｗｘｙｚ０１２３４５６７８９"
_HAN_ALNUM = "ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789"
_ZEN2HAN   = str.maketrans(_ZEN_ALNUM, _HAN_ALNUM)
_ALNUM_PAT = re.compile(r'([A-Za-z0-9Ａ-Ｚａ-ｚ０-９]+)')


def _checker_apply_rules(target_text, rules, for_reporting=False):
    """
    テキストに統一ルールを適用。
    戻り値: [(orig, curr, is_fixed, is_alnum), ...]
    """
    keep_words = list(_CHECKER_KEEP_WORDS)
    for k, v in rules.items():
        if k == v and k not in keep_words:
            keep_words.append(k)
    keep_words = sorted(keep_words, key=len, reverse=True)

    # 保護フレーズをプレースホルダー置換
    protected = target_text
    placeholders = {}
    for i, word in enumerate(keep_words):
        if word in protected:
            ph = f"《《保{i:04d}護》》"
            placeholders[ph] = word
            protected = protected.replace(word, ph)

    # ルール適用
    segments = [(protected, protected, False)]
    for wrong, right in rules.items():
        if wrong == right or not wrong:
            continue
        new_seg = []
        for orig, curr, already in segments:
            if already or str(wrong) not in curr:
                new_seg.append((orig, curr, already))
                continue
            parts = curr.split(str(wrong))
            for j, p in enumerate(parts):
                if p:
                    new_seg.append((p, p, False))
                if j < len(parts) - 1:
                    new_seg.append((str(wrong), str(right), True))
        segments = new_seg

    # プレースホルダーを元に戻す
    restored = []
    for orig, curr, is_fixed in segments:
        t_orig, t_curr = orig, curr
        if not is_fixed:
            for ph, word in placeholders.items():
                t_orig = t_orig.replace(ph, word)
                t_curr = t_curr.replace(ph, word)
        restored.append((t_orig, t_curr, is_fixed))

    # 全角英数字を半角化
    final = []
    for orig, curr, is_fixed in restored:
        if for_reporting and is_fixed:
            new_curr = _ALNUM_PAT.sub(lambda m: m.group(1).translate(_ZEN2HAN), curr)
            has_alnum = bool(_ALNUM_PAT.search(new_curr))
            final.append((orig, new_curr, True, has_alnum))
        else:
            parts = _ALNUM_PAT.split(curr)
            for i2, part in enumerate(parts):
                if not part:
                    continue
                if i2 % 2 == 1:
                    half = part.translate(_ZEN2HAN)
                    was_conv = (half != part)
                    if for_reporting:
                        final.append((part, half, was_conv, True))
                    else:
                        final.append((orig, half, is_fixed or was_conv, False))
                else:
                    if for_reporting:
                        final.append((part, part, False, False))
                    else:
                        final.append((orig, part, is_fixed, False))
    return final


def _checker_is_word_shaded(para):
    """Word 段落に網掛け（背景色）があるか判定"""
    try:
        pPr = para._p.pPr
        if pPr is not None:
            shd = pPr.find(qn('w:shd'))
            if shd is not None:
                val  = shd.get(qn('w:val'))
                fill = shd.get(qn('w:fill'))
                if val and val != 'clear':
                    return True
                if fill and fill not in ('auto', 'FFFFFF', 'clear'):
                    return True
        parent = para._p.getparent()
        if parent is not None and parent.tag.endswith('tc'):
            tcPr = parent.find(qn('w:tcPr'))
            if tcPr is not None:
                shd = tcPr.find(qn('w:shd'))
                if shd is not None:
                    val  = shd.get(qn('w:val'))
                    fill = shd.get(qn('w:fill'))
                    if val and val != 'clear':
                        return True
                    if fill and fill not in ('auto', 'FFFFFF', 'clear'):
                        return True
    except Exception:
        pass
    return False


def _checker_repair_docx(src_path, rules, rgb, out_path, protect_links=True):
    """Word ファイルに統一ルールを適用（修正箇所を指定色で着色）。
    戻り値: True = 1か所以上修正あり / False = 修正なし"""
    doc = Document(str(src_path))
    W_NS = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
    changed = [False]  # 内部関数から書き換えるためリストで保持

    def _process_paragraphs(paragraphs):
        for para in paragraphs:
            is_shaded = _checker_is_word_shaded(para)
            orig_bold = orig_size = None
            if para.runs and para.runs[0].font:
                orig_bold = para.runs[0].font.bold
                orig_size = para.runs[0].font.size

            elements = []
            current_text = ""
            try:
                for child in list(para._p):
                    if child.tag.endswith('hyperlink'):
                        if protect_links:
                            # リンクをカプセルごと退避（テキストは修正しない）
                            if current_text:
                                elements.append({"type": "text", "content": current_text})
                                current_text = ""
                            elements.append({"type": "link", "element": child})
                            para._p.remove(child)
                        else:
                            # リンクのテキストも修正対象にする
                            t_nodes = child.xpath('.//w:t', namespaces={'w': W_NS})
                            text = "".join(t.text for t in t_nodes if t.text)
                            current_text += text
                            para._p.remove(child)
                    elif child.tag.endswith('r') or child.tag.endswith('ins'):
                        t_nodes = child.xpath('.//w:t', namespaces={'w': W_NS})
                        text = "".join(t.text for t in t_nodes if t.text)
                        current_text += text
                        para._p.remove(child)
                if current_text:
                    elements.append({"type": "text", "content": current_text})
            except Exception:
                elements = [{"type": "text", "content": para.text}]
                para.text = ""

            for el in elements:
                if el["type"] == "text":
                    if not el["content"]:
                        continue
                    parts = _checker_apply_rules(el["content"], rules, for_reporting=False)
                    for orig, curr, is_fixed, is_alnum in parts:
                        run = para.add_run(curr)
                        run.font.name = 'ＭＳ 明朝'
                        run._element.rPr.rFonts.set(qn('w:eastAsia'), 'ＭＳ 明朝')
                        if is_shaded:
                            if orig_size is not None:
                                run.font.size = orig_size
                            if orig_bold is not None:
                                run.font.bold = orig_bold
                        else:
                            run.font.size = DocxPt(10.5)
                        if is_fixed:
                            changed[0] = True
                            run.font.color.rgb = DocxRGBColor(rgb[0], rgb[1], rgb[2])
                            run.bold = False
                elif el["type"] == "link":
                    para._p.append(el["element"])

    _process_paragraphs(doc.paragraphs)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                _process_paragraphs(cell.paragraphs)
    if changed[0]:
        doc.save(str(out_path))
    return changed[0]


def _checker_repair_xlsx(src_path, rules, rgb, out_path, protect_links=True):
    """Excel ファイルに統一ルールを適用（修正箇所を指定色で着色）。
    戻り値: True = 1か所以上修正あり / False = 修正なし"""
    wb = openpyxl.load_workbook(str(src_path))
    hex_color = f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    changed = False
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if not (cell.value and isinstance(cell.value, str)):
                    continue
                if protect_links and cell.hyperlink:
                    continue
                is_shaded = False
                if (cell.fill and cell.fill.patternType and
                        cell.fill.patternType != 'none'):
                    fc = cell.fill.fgColor.rgb
                    if fc and fc not in ('00000000', 'FFFFFFFF', '00FFFFFF'):
                        is_shaded = True
                orig_bold = cell.font.bold if cell.font else False
                orig_size = cell.font.size if cell.font else 11
                parts = _checker_apply_rules(cell.value, rules, for_reporting=False)
                if any(p[2] for p in parts) or any(p[3] for p in parts) or not is_shaded:
                    cell.value = "".join(p[1] for p in parts)
                    is_fixed_present = any(p[2] for p in parts)
                    if is_fixed_present:
                        changed = True
                    new_color = hex_color if is_fixed_present else (
                        cell.font.color if cell.font else None)
                    new_size = orig_size if is_shaded else 10.5
                    new_bold = orig_bold if is_shaded else False
                    cell.font = OpenpyxlFont(
                        name='ＭＳ 明朝', size=new_size,
                        color=new_color, bold=new_bold)
    if changed:
        wb.save(str(out_path))
    return changed


def _checker_repair_pptx(src_path, rules, rgb, out_path, protect_links=True):
    """PowerPoint ファイルに統一ルールを適用（修正箇所を指定色で着色）。
    戻り値: True = 1か所以上修正あり / False = 修正なし"""
    prs = PptxPresentation(str(src_path))
    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            is_shaded = hasattr(shape, "fill") and shape.fill.type == 1
            if not (hasattr(shape, "text_frame") and shape.text_frame):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if protect_links:
                    has_link = any(
                        hasattr(run, "hyperlink") and run.hyperlink and run.hyperlink.address
                        for run in paragraph.runs
                    )
                    if has_link:
                        continue
                combined = "".join(run.text for run in paragraph.runs)
                parts = _checker_apply_rules(combined, rules, for_reporting=False)
                if any(p[2] for p in parts) or any(p[3] for p in parts) or not is_shaded:
                    paragraph.text = ""
                    for orig, curr, is_fixed, is_alnum in parts:
                        new_run = paragraph.add_run()
                        new_run.text = curr
                        new_run.font.name = 'ＭＳ 明朝'
                        if not is_shaded:
                            new_run.font.size = PptxPt(10.5)
                        if is_fixed:
                            changed = True
                            new_run.font.color.rgb = PptxRGBColor(rgb[0], rgb[1], rgb[2])
                            new_run.font.bold = False
    if changed:
        prs.save(str(out_path))
    return changed


def _checker_check_pdf(src_path, rules):
    """PDF をチェックして修正推奨リストを返す"""
    results = []
    invisible = re.compile(r'[\s\u200B-\u200F\u202A-\u202E\u2060-\u206F\uFEFF\u00A0]+')
    with pdfplumber.open(str(src_path)) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text(x_tolerance=2, y_tolerance=2)
            if not text:
                continue
            pure = re.sub(invisible, '', text)
            parts = _checker_apply_rules(pure, rules, for_reporting=True)
            full_text = "".join(p[1] for p in parts)
            idx = 0
            for orig, curr, is_fixed, is_alnum in parts:
                if is_fixed:
                    s = max(0, idx - 15)
                    e = min(len(full_text), idx + len(curr) + 15)
                    ctx = full_text[s:e]
                    reason = ("英数字の半角化"
                              if orig.translate(_ZEN2HAN) == curr and orig != curr
                              else "統一ルールの適用")
                    results.append({
                        "ページ":    i + 1,
                        "NGワード":  orig,
                        "修正案":    curr,
                        "修正理由":  reason,
                        "周辺の文章": f"…{ctx}…",
                    })
                idx += len(curr)
    return results


def _checker_expand_paths(input_paths, exts):
    """
    パスリストを展開する。
    - ファイル → そのまま追加（対象外拡張子はスキップ）
    - フォルダ → 再帰的に対象ファイルを収集（階層制限なし）
    戻り値: Path のリスト（重複除去・ソート済み）
    """
    seen   = set()
    result = []
    for raw in input_paths:
        p = Path(raw)
        if p.is_dir():
            for ext in exts:
                for found in sorted(p.rglob(f"*{ext}")):
                    key = found.resolve()
                    if key not in seen:
                        seen.add(key)
                        result.append(found)
        elif p.is_file():
            if p.suffix.lower() in exts:
                key = p.resolve()
                if key not in seen:
                    seen.add(key)
                    result.append(p)
    return result


def _checker_write_pdf_report(pdf_reports, out_path):
    """
    PDF チェック結果を 1 つの Word 文書にまとめて出力する。
    pdf_reports: [(filename, [{"ページ", "NGワード", "修正案", "修正理由", "周辺の文章"}, ...]), ...]
    """
    from datetime import datetime
    from docx.oxml.ns import qn as _qn
    from docx.oxml   import OxmlElement as _Elem

    doc = Document()

    # ── 既定フォントを游明朝に統一 ───────────────────────────────
    style = doc.styles['Normal']
    style.font.name = 'ＭＳ 明朝'
    style.font.size = DocxPt(10.5)

    # ── タイトル ──────────────────────────────────────────────────
    title_p = doc.add_heading('PDF 統一語句チェック結果', level=0)
    title_p.alignment = 1  # CENTER
    for run in title_p.runs:
        run.font.name = 'ＭＳ 明朝'

    date_str = datetime.now().strftime('%Y年%m月%d日  %H:%M')
    info_p = doc.add_paragraph(f'作成日時: {date_str}')
    info_p.alignment = 1
    doc.add_paragraph()

    # ── 各 PDF のセクション ───────────────────────────────────────
    for fname, report in pdf_reports:
        # ファイル名見出し（H1）
        h = doc.add_heading(fname, level=1)
        for run in h.runs:
            run.font.name = 'ＭＳ 明朝'

        if not report:
            ok_p = doc.add_paragraph('　✓  修正箇所なし（統一ルールに準拠しています）')
            if ok_p.runs:
                ok_p.runs[0].font.color.rgb = DocxRGBColor(0, 128, 0)
            doc.add_paragraph()
            continue

        # 指摘件数
        cnt_p = doc.add_paragraph(f'　⚠  {len(report)} 箇所に修正を推奨します')
        if cnt_p.runs:
            cnt_p.runs[0].font.color.rgb = DocxRGBColor(200, 80, 0)

        # テーブル
        table = doc.add_table(rows=1, cols=5)
        table.style = 'Table Grid'
        tbl_hdr = table.rows[0].cells
        for i, hdr_txt in enumerate(
                ['ページ', 'NGワード', '修正案', '修正理由', '周辺の文章']):
            tbl_hdr[i].text = hdr_txt
            run = tbl_hdr[i].paragraphs[0].runs[0]
            run.font.bold = True
            run.font.name = 'ＭＳ 明朝'
            # ヘッダー背景色（薄いグレー）
            tc_pr = tbl_hdr[i]._tc.get_or_add_tcPr()
            shd   = _Elem('w:shd')
            shd.set(_qn('w:val'),  'clear')
            shd.set(_qn('w:color'), 'auto')
            shd.set(_qn('w:fill'), 'D9D9D9')
            tc_pr.append(shd)

        for r in report:
            row = table.add_row().cells
            row[0].text = str(r['ページ'])
            row[1].text = r['NGワード']
            row[2].text = r['修正案']
            row[3].text = r['修正理由']
            row[4].text = r['周辺の文章']
            for cell in row:
                for para in cell.paragraphs:
                    for run in para.runs:
                        run.font.name = 'ＭＳ 明朝'
                        run.font.size = DocxPt(9)

        doc.add_paragraph()

    doc.save(str(out_path))


# ════════════════════════════════════════════════════════════════
#  統一語句ルール編集ダイアログ
# ════════════════════════════════════════════════════════════════

class RulesEditorDialog(tk.Toplevel):
    """統一語句ルールを GUIで追加・編集・削除して CSV へ保存するダイアログ"""

    def __init__(self, parent, rules: dict, save_path: Path, on_saved):
        super().__init__(parent)
        self.title("統一語句ルール編集")
        self.configure(bg=C["bg"])
        self.geometry("640x540")
        self.minsize(500, 400)
        self.resizable(True, True)
        self._save_path = save_path
        self._on_saved  = on_saved
        self._row_data  = []   # list of (wrong_var, right_var, frame)
        self._build()
        self._load_rules(rules)
        self.grab_set()         # モーダル

    # ── UI構築 ────────────────────────────────────────────────────

    def _build(self):
        # ヘッダー
        hdr = tk.Frame(self, bg=C["surface"], padx=16, pady=10)
        hdr.pack(fill=tk.X)
        tk.Label(hdr, text="⚓  統一語句ルール編集",
                 font=("Helvetica", 13, "bold"),
                 bg=C["surface"], fg=C["text"]).pack(side=tk.LEFT)
        tk.Label(hdr,
                 text="（NGワード → 統一語句）を自由に追加・修正・削除できます",
                 font=("Helvetica", 9), bg=C["surface"], fg=C["sub"]
                 ).pack(side=tk.LEFT, padx=12)

        # 列ヘッダー
        col_hdr = tk.Frame(self, bg=C["accent"])
        col_hdr.pack(fill=tk.X, padx=16, pady=(10, 0))
        tk.Label(col_hdr, text="  削除", width=5,
                 font=("Helvetica", 8, "bold"),
                 bg=C["accent"], fg=C["sub"]).grid(row=0, column=0, padx=(4, 2), pady=4)
        tk.Label(col_hdr, text="NGワード（修正前）",
                 font=("Helvetica", 9, "bold"),
                 bg=C["accent"], fg=C["text"]).grid(row=0, column=1, padx=4, pady=4, sticky="w")
        tk.Label(col_hdr, text="→",
                 font=("Helvetica", 9), bg=C["accent"], fg=C["sub"]
                 ).grid(row=0, column=2, padx=2)
        tk.Label(col_hdr, text="統一語句（修正後）",
                 font=("Helvetica", 9, "bold"),
                 bg=C["accent"], fg=C["text"]).grid(row=0, column=3, padx=4, pady=4, sticky="w")
        col_hdr.columnconfigure(1, weight=1)
        col_hdr.columnconfigure(3, weight=1)

        # スクロール可能な行エリア
        scroll_wrap = tk.Frame(self, bg=C["bg"])
        scroll_wrap.pack(fill=tk.BOTH, expand=True, padx=16, pady=(0, 0))
        canvas = tk.Canvas(scroll_wrap, bg=C["bg"], highlightthickness=0)
        sb = ttk.Scrollbar(scroll_wrap, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=sb.set)
        sb.pack(side=tk.RIGHT, fill=tk.Y)
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        self._rows_frame = tk.Frame(canvas, bg=C["bg"])
        fid = canvas.create_window((0, 0), window=self._rows_frame, anchor="nw")
        self._rows_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.bind("<Configure>",
                    lambda e: canvas.itemconfig(fid, width=e.width))
        # マウスホイールスクロール
        def _on_wheel(e):
            canvas.yview_scroll(int(-1 * (e.delta / 120)), "units")
        canvas.bind_all("<MouseWheel>", _on_wheel)
        self._canvas = canvas

        # ボタン行
        btn_row = tk.Frame(self, bg=C["bg"])
        btn_row.pack(fill=tk.X, padx=16, pady=10)
        FlatButton(btn_row, text="＋ 行を追加",
                   command=lambda: self._add_row("", ""),
                   bg=C["accent"], fg=C["text"],
                   font=("Helvetica", 9)
                   ).pack(side=tk.LEFT)
        FlatButton(btn_row, text="✕ 閉じる",
                   command=self.destroy,
                   bg=C["accent"], fg=C["text"],
                   font=("Helvetica", 9)
                   ).pack(side=tk.RIGHT, padx=(8, 0))
        FlatButton(btn_row, text="💾 保存して閉じる",
                   command=self._save,
                   bg=C["ok"], fg=C["bg"],
                   font=("Helvetica", 10), bold=True
                   ).pack(side=tk.RIGHT)

    # ── データ操作 ────────────────────────────────────────────────

    def _load_rules(self, rules: dict):
        for wrong, right in rules.items():
            self._add_row(wrong, right)

    def _add_row(self, wrong: str = "", right: str = ""):
        rf = tk.Frame(self._rows_frame, bg=C["surface"], pady=2)
        rf.pack(fill=tk.X, pady=1, padx=2)

        wrong_var = tk.StringVar(value=wrong)
        right_var = tk.StringVar(value=right)

        # 削除ボタン
        def _delete(frame=rf, wv=wrong_var, rv=right_var):
            frame.destroy()
            self._row_data = [(w, r, f) for w, r, f in self._row_data
                              if f is not frame]

        FlatButton(rf, text="✕", command=_delete,
                   bg=C["err"], fg=C["text"],
                   font=("Helvetica", 8), padx=5, pady=2
                   ).pack(side=tk.LEFT, padx=(4, 4))

        # NGワード 入力
        tk.Entry(rf, textvariable=wrong_var, width=20,
                 bg=C["input_bg"], fg=C["text"],
                 insertbackground=C["text"],
                 relief=tk.FLAT, font=("Helvetica", 10)
                 ).pack(side=tk.LEFT, padx=(0, 4), ipady=4, fill=tk.X, expand=True)

        tk.Label(rf, text="→", bg=C["surface"], fg=C["sub"],
                 font=("Helvetica", 10)).pack(side=tk.LEFT)

        # 統一語句 入力
        tk.Entry(rf, textvariable=right_var, width=20,
                 bg=C["input_bg"], fg=C["text"],
                 insertbackground=C["text"],
                 relief=tk.FLAT, font=("Helvetica", 10)
                 ).pack(side=tk.LEFT, padx=(4, 8), ipady=4, fill=tk.X, expand=True)

        self._row_data.append((wrong_var, right_var, rf))

    def _save(self):
        # 生きている行だけ収集
        rows = []
        for wrong_var, right_var, rf in self._row_data:
            try:
                if not rf.winfo_exists():
                    continue
            except tk.TclError:
                continue
            w = wrong_var.get().strip()
            r = right_var.get().strip()
            if w:
                rows.append((w, r or w))

        if not rows:
            messagebox.showwarning("空のルール",
                                   "1件以上ルールを入力してください。",
                                   parent=self)
            return

        try:
            with open(self._save_path, "w", encoding="utf-8-sig", newline="") as f:
                writer = csv.writer(f)
                writer.writerow(["類義語", "統一語句"])
                for w, r in rows:
                    writer.writerow([w, r])
            messagebox.showinfo("保存完了",
                                f"{len(rows)} 件のルールを保存しました。\n{self._save_path.name}",
                                parent=self)
            self._on_saved(self._save_path)
            self.destroy()
        except Exception as e:
            messagebox.showerror("保存エラー", str(e), parent=self)


# ════════════════════════════════════════════════════════════════
#  統一ルール修正チェッカー 画面
# ════════════════════════════════════════════════════════════════

class CheckerFrame(LoggedFrame):

    COLOR_MAP = {
        "赤": (255, 0, 0),
        "青": (0, 0, 255),
        "緑": (0, 128, 0),
        "黒": (0, 0, 0),
    }
    CHECKER_EXTS = {".docx", ".xlsx", ".pptx", ".pdf"}

    def __init__(self, parent, on_back, on_go_pdf, on_go_link):
        super().__init__(parent, bg=C["bg"])
        self.on_back    = on_back
        self.on_go_pdf  = on_go_pdf
        self.on_go_link = on_go_link
        self._cancel_flag  = threading.Event()
        self._rules        = {}
        self._rules_path   = None   # 現在読み込み中の CSV パス
        self._build()
        self._try_load_default_rules()

    def _build(self):
        build_header(self, "⚓  統一ルール修正チェッカー",
                     "Word / Excel / PowerPoint / PDF の統一ルール適用")
        nav = tk.Frame(self, bg=C["bg"])
        nav.pack(fill=tk.X, padx=16, pady=(8, 0))
        nav_button(nav, "← ホームへ戻る",   self.on_back).pack(side=tk.LEFT)
        nav_button(nav, "PDF変換へ →",       self.on_go_pdf).pack(side=tk.RIGHT)
        nav_button(nav, "リンク設定へ →",    self.on_go_link).pack(side=tk.RIGHT, padx=(0, 6))

        _, main = make_scrollable_frame(self)
        pad = dict(padx=16, pady=(0, 10))

        # ── 入力ファイル DropZone ─────────────────────────────────
        self.src_zone = DropZone(
            main, "チェック対象ファイル",
            "Word(.docx) / Excel(.xlsx) / PowerPoint(.pptx) / PDF をドロップ",
            select_mode="file", allow_multiple=True,
            count_extensions=self.CHECKER_EXTS
        )
        self.src_zone.pack(fill=tk.X, **pad)

        # ── 出力先 DropZone ───────────────────────────────────────
        self.out_zone = DropZone(
            main, "出力先フォルダ（省略時：元ファイルと同じ場所）",
            "フォルダをドロップ、または空のまま実行",
            select_mode="folder", allow_multiple=False,
            count_extensions=set()
        )
        self.out_zone.pack(fill=tk.X, **pad)

        # ── ルールCSV 選択・編集行 ────────────────────────────────
        section_divider(main)
        rules_row = tk.Frame(main, bg=C["bg"])
        rules_row.pack(fill=tk.X, padx=16, pady=(0, 8))
        tk.Label(rules_row, text="ルール CSV :",
                 font=("Helvetica", 10), bg=C["bg"], fg=C["sub"]
                 ).pack(side=tk.LEFT)
        self.rules_label = tk.Label(
            rules_row, text="（読み込み中…）",
            font=("Helvetica", 9), bg=C["bg"], fg=C["info"],
            wraplength=300, anchor="w"
        )
        self.rules_label.pack(side=tk.LEFT, padx=8, fill=tk.X, expand=True)
        FlatButton(rules_row, text="✏ 編集",
                   command=self._open_rules_editor,
                   bg=C["ok"], fg=C["bg"],
                   font=("Helvetica", 9), bold=True
                   ).pack(side=tk.RIGHT)
        FlatButton(rules_row, text="CSV を選択",
                   command=self._pick_rules_csv,
                   bg=C["accent"], fg=C["text"],
                   font=("Helvetica", 9)
                   ).pack(side=tk.RIGHT, padx=(0, 6))
        FlatButton(rules_row, text="リセット",
                   command=self._try_load_default_rules,
                   bg=C["accent"], fg=C["text"],
                   font=("Helvetica", 9)
                   ).pack(side=tk.RIGHT, padx=(0, 6))

        # ── 修正箇所の色 ──────────────────────────────────────────
        color_row = tk.Frame(main, bg=C["bg"])
        color_row.pack(fill=tk.X, padx=16, pady=(0, 6))
        tk.Label(color_row, text="修正箇所の色 :",
                 font=("Helvetica", 10), bg=C["bg"], fg=C["sub"]
                 ).pack(side=tk.LEFT)
        self._color_var = tk.StringVar(value="赤")
        for color_name in self.COLOR_MAP:
            tk.Radiobutton(
                color_row, text=color_name,
                variable=self._color_var, value=color_name,
                bg=C["bg"], fg=C["text"],
                selectcolor=C["accent"],
                activebackground=C["bg"], activeforeground=C["text"],
                font=("Helvetica", 10)
            ).pack(side=tk.LEFT, padx=(8, 0))

        # ── ハイパーリンク保護オプション ──────────────────────────
        link_row = tk.Frame(main, bg=C["bg"])
        link_row.pack(fill=tk.X, padx=16, pady=(0, 12))
        self._protect_links_var = tk.BooleanVar(value=True)
        tk.Checkbutton(
            link_row,
            text="ハイパーリンクを保護する（リンク文字列は修正しない）",
            variable=self._protect_links_var,
            bg=C["bg"], fg=C["sub"],
            selectcolor=C["accent"],
            activebackground=C["bg"], activeforeground=C["text"],
            font=("Helvetica", 9)
        ).pack(side=tk.LEFT)

        # ── 実行・キャンセルボタン ────────────────────────────────
        section_divider(main)
        btn_row = tk.Frame(main, bg=C["bg"])
        btn_row.pack(fill=tk.X, padx=16, pady=(0, 10))
        self.run_btn = FlatButton(
            btn_row, text="▶  修正チェックを開始",
            command=self._start,
            bg=C["primary"], fg=C["text"],
            font=("Helvetica", 11), bold=True
        )
        self.run_btn.pack(side=tk.LEFT, fill=tk.X, expand=True)
        self.cancel_btn = FlatButton(
            btn_row, text="■  キャンセル",
            command=self._cancel,
            bg=C["accent"], fg=C["text"],
            font=("Helvetica", 9), state=tk.DISABLED
        )
        self.cancel_btn.pack(side=tk.RIGHT, padx=(8, 0))

        # ── ログ ──────────────────────────────────────────────────
        self.log_txt = make_log_widget(main)

    # ── rules.csv ────────────────────────────────────────────────

    def _try_load_default_rules(self):
        """スクリプトと同じフォルダの rules.csv を自動読み込み。
        PyInstaller バンドル時は sys._MEIPASS を参照する。"""
        if getattr(sys, 'frozen', False):
            base = Path(sys._MEIPASS)   # PyInstaller 展開先
        else:
            base = Path(__file__).parent
        default = base / "rules.csv"
        if default.exists():
            self._load_rules_from(default)
        else:
            self._rules      = {}
            self._rules_path = None
            self.rules_label.configure(
                text="（rules.csv 未検出 — CSV を手動で選択してください）",
                fg=C["warn"])

    def _pick_rules_csv(self):
        path = filedialog.askopenfilename(
            title="ルール CSV を選択",
            filetypes=[("CSV ファイル", "*.csv"), ("すべて", "*.*")]
        )
        if path:
            self._load_rules_from(Path(path))

    def _load_rules_from(self, path: Path):
        rules = _checker_load_rules(path)
        if rules:
            self._rules      = rules
            self._rules_path = path
            self.rules_label.configure(
                text=f"{path.name}  ({len(rules)} ルール読み込み済み)",
                fg=C["ok"])
            self._log(f"ルール読み込み完了: {path.name}  ({len(rules)} 件)", "success")
        else:
            self._rules      = {}
            self._rules_path = path   # パスは覚えておく（新規編集できるように）
            self.rules_label.configure(
                text=f"読み込み失敗: {path.name}", fg=C["err"])
            self._log(f"ルール CSV の読み込みに失敗しました: {path}", "error")

    def _open_rules_editor(self):
        """統一語句ルール編集ダイアログを開く"""
        # 保存先パスを決定（未選択なら rules.csv をデフォルト）
        save_path = self._rules_path or (Path(__file__).parent / "rules.csv")
        RulesEditorDialog(
            parent=self,
            rules=self._rules,
            save_path=save_path,
            on_saved=lambda p: self._load_rules_from(p)
        )

    # ── 処理制御 ──────────────────────────────────────────────────

    def _start(self):
        paths = self.src_zone.selected_paths
        if not paths:
            messagebox.showwarning("ファイル未指定",
                                   "チェック対象ファイルを指定してください。")
            return
        if not self._rules:
            messagebox.showwarning("ルール未設定",
                                   "ルール CSV が読み込まれていません。\n"
                                   "CSV を選択してください。")
            return
        out_paths     = self.out_zone.selected_paths
        out_dir       = Path(out_paths[0]) if out_paths else None
        rgb           = self.COLOR_MAP[self._color_var.get()]
        protect_links = self._protect_links_var.get()
        self._cancel_flag.clear()
        self.run_btn.configure(text="処理中…", state=tk.DISABLED)
        self.cancel_btn.configure(state=tk.NORMAL)
        threading.Thread(
            target=self._run_worker,
            args=(paths, out_dir, rgb, protect_links),
            daemon=True
        ).start()

    def _cancel(self):
        self._cancel_flag.set()
        self._log("キャンセルを要求しました…", "warning")

    def _run_worker(self, paths, out_dir, rgb, protect_links):
        ok = err = skip = 0
        pdf_reports = []       # [(filename, report_list), ...]
        all_pdf_srcs = []      # PDF ソースパス（レポート保存先フォールバック用）

        try:
            # ── 入力パスをグループ化 ─────────────────────────────────
            # groups: [ (root_or_None, [Path, ...]) ]
            #   root_or_None … D&D したフォルダ、または None（単体ファイル）
            groups = []
            seen   = set()
            for raw in paths:
                p = Path(raw)
                if p.is_dir():
                    group_files = []
                    for ext in self.CHECKER_EXTS:
                        for found in sorted(p.rglob(f"*{ext}")):
                            key = found.resolve()
                            if key not in seen:
                                seen.add(key)
                                group_files.append(found)
                    group_files.sort()
                    if group_files:
                        groups.append((p, group_files))
                elif p.is_file():
                    if p.suffix.lower() in self.CHECKER_EXTS:
                        key = p.resolve()
                        if key not in seen:
                            seen.add(key)
                            groups.append((None, [p]))

            total_files = sum(len(g[1]) for g in groups)
            if not total_files:
                self._log("対象ファイルが見つかりませんでした。", "warning")
                return

            self._log(f"処理対象: {total_files} ファイル", "info")

            # ── グループ単位で処理 ──────────────────────────────────
            for root, files in groups:
                if self._cancel_flag.is_set():
                    self._log("⚠ キャンセルされました。", "warning")
                    break

                is_folder = (root is not None)
                if is_folder:
                    # フォルダグループの出力ベース
                    base_out = out_dir if out_dir else (
                        root.parent / (root.name + "_チェック済"))
                    self._log(f"── フォルダ: {root.name} ({len(files)} ファイル)", "info")

                # group_results: [(src, dest, result)]
                #   result: True=修正あり, False=修正なし, None=PDF, "error"=エラー
                group_results    = []
                group_has_change = False   # Office ファイルに修正があったか

                for src in files:
                    if self._cancel_flag.is_set():
                        break

                    ext = src.suffix.lower()

                    # ── 出力先パスの決定 ──────────────────────────
                    if is_folder:
                        # フォルダ階層を保持: base_out / フォルダ名 / 相対パス
                        rel  = src.relative_to(root)
                        dest = base_out / root.name / rel
                    elif out_dir:
                        dest = out_dir / src.name
                    else:
                        dest = src.parent / src.name

                    # 単体ファイルの同一パスコンフリクト回避
                    if not is_folder and dest.resolve() == src.resolve():
                        counter = 1
                        while dest.resolve() == src.resolve() or dest.exists():
                            dest = dest.parent / f"{src.stem}({counter}){src.suffix}"
                            counter += 1

                    # ── ファイル種別ごとの処理 ───────────────────
                    if ext == ".pdf":
                        if not _PDFPLUMBER_OK:
                            self._log(
                                "  ✗ pdfplumber 未インストール（pip install pdfplumber）",
                                "error")
                            group_results.append((src, dest, "error"))
                            continue
                        self._log(f"PDF チェック中: {src.name}", "info")
                        try:
                            report = _checker_check_pdf(src, self._rules)
                            pdf_reports.append((src.name, report))
                            all_pdf_srcs.append(src)
                            if report:
                                self._log(f"  ⚠ {len(report)} 箇所で修正を推奨", "warning")
                            else:
                                self._log("  ✓ 問題なし", "success")
                            group_results.append((src, dest, None))   # None = PDF
                        except Exception as e:
                            self._log(f"  ✗ エラー: {e}", "error")
                            group_results.append((src, dest, "error"))

                    else:
                        # Word / Excel / PPT
                        self._log(f"処理中: {src.name}", "info")
                        try:
                            dest.parent.mkdir(parents=True, exist_ok=True)
                            if ext == ".docx":
                                modified = _checker_repair_docx(
                                    src, self._rules, rgb, dest,
                                    protect_links=protect_links)
                            elif ext == ".xlsx":
                                if not _OPENPYXL_OK:
                                    raise RuntimeError(
                                        "openpyxl 未インストール（pip install openpyxl）")
                                modified = _checker_repair_xlsx(
                                    src, self._rules, rgb, dest,
                                    protect_links=protect_links)
                            elif ext == ".pptx":
                                if not _PPTX_OK:
                                    raise RuntimeError(
                                        "python-pptx 未インストール（pip install python-pptx）")
                                modified = _checker_repair_pptx(
                                    src, self._rules, rgb, dest,
                                    protect_links=protect_links)
                            else:
                                modified = False
                            group_results.append((src, dest, modified))
                            if modified:
                                group_has_change = True
                        except Exception as e:
                            self._log(f"  ✗ エラー: {e}", "error")
                            group_results.append((src, dest, "error"))

                # ── グループ後処理（カウント & ファイル出力） ──────
                if is_folder and not group_has_change:
                    # フォルダ内に Office 修正なし → 全スキップ
                    for src, dest, result in group_results:
                        if result == "error":
                            err += 1
                        else:
                            # 修正ありとして保存されたファイルがあれば削除
                            if dest.exists():
                                try:
                                    dest.unlink()
                                except Exception:
                                    pass
                            skip += 1
                    self._log(
                        f"  フォルダ「{root.name}」: 修正箇所なし"
                        f"（{len(group_results)} ファイルをスキップ）",
                        "success")
                else:
                    for src, dest, result in group_results:
                        if result == "error":
                            err += 1

                        elif result is None:
                            # PDF ファイル
                            if is_folder and group_has_change:
                                # フォルダに修正あり → PDF も出力
                                dest.parent.mkdir(parents=True, exist_ok=True)
                                shutil.copy2(str(src), str(dest))
                                self._log(
                                    f"  ✓ 出力（PDF）: {src.name}", "success")
                            ok += 1   # チェック済みとしてカウント

                        elif result is True:
                            # 修正あり
                            self._log(
                                f"  ✓ 出力（修正あり）: {src.name}", "success")
                            ok += 1

                        else:
                            # result is False → 修正なし
                            if dest.exists():
                                try:
                                    dest.unlink()
                                except Exception:
                                    pass
                            if is_folder:
                                # フォルダに修正あり → 元ファイルを出力
                                dest.parent.mkdir(parents=True, exist_ok=True)
                                shutil.copy2(str(src), str(dest))
                                self._log(
                                    f"  ✓ 出力（修正なし）: {src.name}", "success")
                                ok += 1
                            else:
                                # 単体ファイル → スキップ
                                self._log(
                                    f"  ✓ 修正箇所なし（出力スキップ）: {src.name}",
                                    "success")
                                skip += 1

            # ── PDF 結果を Word にまとめて出力 ────────────────────
            # ※ out_dir 直下に単体で保存（フォルダ階層の外）
            if pdf_reports:
                self._log("─" * 50, "")
                self._log("【PDF チェック結果】", "info")
                total_issues = 0
                for fname, report in pdf_reports:
                    self._log(
                        f"  {fname}: {len(report)} 件",
                        "warning" if report else "success")
                    total_issues += len(report)

                if total_issues == 0:
                    self._log(
                        "  ✓ すべての PDF に修正箇所なし（レポート出力スキップ）",
                        "success")
                else:
                    # Word レポートは out_dir 直下（フォルダ階層の外側）
                    if out_dir:
                        report_dir = out_dir
                    elif all_pdf_srcs:
                        report_dir = all_pdf_srcs[0].parent
                    else:
                        report_dir = Path.cwd()
                    report_dir.mkdir(parents=True, exist_ok=True)

                    report_path = report_dir / "PDF_統一語句チェック結果.docx"
                    counter = 1
                    while report_path.exists():
                        report_path = (
                            report_dir / f"PDF_統一語句チェック結果({counter}).docx")
                        counter += 1

                    try:
                        _checker_write_pdf_report(pdf_reports, report_path)
                        self._log(
                            f"  ✓ Word レポート保存: {report_path.name}  "
                            f"（全 {len(pdf_reports)} PDF / 指摘 {total_issues} 件）",
                            "success")
                    except Exception as e:
                        self._log(f"  ✗ Word レポート保存エラー: {e}", "error")

                self._log("─" * 50, "")

            self._log(
                f"完了: {ok} 件成功  {err} 件エラー  {skip} 件スキップ",
                "success" if err == 0 else "warning")
        except Exception as e:
            self._log(f"予期しないエラー: {e}", "error")
        finally:
            self.after(0, self._reset_btns)

    def _reset_btns(self):
        self.cancel_btn.configure(state=tk.DISABLED)
        self.run_btn.configure(text="▶  修正チェックを開始", state=tk.NORMAL)

    def _log(self, msg, tag=""):
        self.after(0, lambda m=msg, t=tag: log_write(self.log_txt, m, t))


# ════════════════════════════════════════════════════════════════
#  メインアプリ（画面切り替えコントローラ）
# ════════════════════════════════════════════════════════════════

class RakurakuJC:

    def __init__(self):
        if _DND_BACKEND == "tkdnd":
            try:
                self.root = TkinterDnD.Tk()
            except Exception:
                self.root = tk.Tk()
        else:
            self.root = tk.Tk()

        self.root.title(f"楽々JC  v{APP_VERSION}")
        self.root.geometry("680x800")
        self.root.minsize(560, 680)
        self.root.configure(bg=C["bg"])

        self._current = None
        self._show_launcher()

        # 描画を確定させてからウィンドウを前面に出す
        self.root.update_idletasks()
        self.root.update()

        if platform.system() == "Darwin":
            self.root.lift()
            self.root.attributes("-topmost", True)
            self.root.after(200, lambda: self.root.attributes("-topmost", False))
            self.root.after(300, lambda: self.root.focus_force())

    def check_update(self):
        def _do():
            try:
                url = GITHUB_RAW + "/version.txt"
                with urllib.request.urlopen(url, timeout=5) as res:
                    latest = res.read().decode().strip()
                if latest == APP_VERSION:
                    self.root.after(0, lambda: messagebox.showinfo(
                        "アップデート確認", "最新バージョンです！\n現在: v" + APP_VERSION))
                else:
                    lv = latest
                    self.root.after(0, lambda: messagebox.showinfo(
                        "アップデートあり",
                        f"新しいバージョンがあります。\n現在: v{APP_VERSION}  最新: v{lv}"))
            except Exception as e:
                self.root.after(0, lambda: messagebox.showerror("エラー", str(e)))
        threading.Thread(target=_do, daemon=True).start()

    def _clear(self):
        if self._current:
            self._current.destroy()
            self._current = None

    def _show_launcher(self):
        self._clear()
        f = LauncherFrame(
            self.root,
            on_link=self._show_link,
            on_pdf=self._show_pdf,
            on_checker=self._show_checker
        )
        f.pack(fill=tk.BOTH, expand=True)
        self._current = f

    def _show_link(self):
        self._clear()
        f = LinkFrame(
            self.root,
            on_back=self._show_launcher,
            on_go_pdf=self._show_pdf,
            on_go_checker=self._show_checker,
            dnd_ok=(_DND_BACKEND == "tkdnd"),
            on_check_update=self.check_update
        )
        f.pack(fill=tk.BOTH, expand=True)
        self._current = f

    def _show_pdf(self):
        self._clear()
        f = PdfFrame(
            self.root,
            on_back=self._show_launcher,
            on_go_link=self._show_link,
            on_go_checker=self._show_checker
        )
        f.pack(fill=tk.BOTH, expand=True)
        self._current = f

    def _show_checker(self):
        self._clear()
        f = CheckerFrame(
            self.root,
            on_back=self._show_launcher,
            on_go_pdf=self._show_pdf,
            on_go_link=self._show_link
        )
        f.pack(fill=tk.BOTH, expand=True)
        self._current = f

    def run(self):
        self.root.mainloop()


# ════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    RakurakuJC().run()
